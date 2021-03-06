import os,zipfile,re,operator,pptx
from django.shortcuts import render
from django.http import HttpResponseRedirect, JsonResponse
from django.core.urlresolvers import reverse
from django.db.models import TextField, CharField, Q, Count
from django.db.models.functions import Lower
from django.contrib.auth.decorators import login_required
from itertools import chain
from django.conf import settings
from django.core.files import File
from pptx import Presentation
from subprocess import call

from haystack.query import SearchQuerySet

# path to media root for adding a zip archive
MEDIA_ROOT = settings.MEDIA_ROOT

# import all website models and forms
from website.models import *
from website.forms import *

#for email
from django.core.mail import send_mail

#for csv export
from djqscsv import render_to_csv_response

# GAHTC Views
def index(request):


	"""
	  Index page
	"""
	context_dict = {}
	return render(request, 'website/index.html', context_dict)

def about(request):


	"""
	  About page
	"""
	context_dict = {}
	return render(request, 'website/about.html', context_dict)

def howToUse(request):


	"""
	  How to Use GAHTC.org page
	"""
	context_dict = {}
	return render(request, 'website/how_to_use.html', context_dict)

def timescape(request):


	"""
	  timescape page
	"""
	context_dict = {}
	return render(request, 'website/timescape.html', context_dict)

def howtobecomeamember(request):
	context_dict = {}
	return render(request, 'website/howtobecomeamember.html', context_dict)

def membersconference(request):
	context_dict = {}
	return render(request, 'website/membersconference.html', context_dict)

def testimonials(request):
	context_dict = {}
	return render(request, 'website/testimonials.html', context_dict)

def grants(request):
	context_dict = {}
	return render(request, 'website/grants.html', context_dict)


def contact(request):


	"""
	  contact page
	"""
	context_dict = {}
	return render(request, 'website/contact.html', context_dict)

def architectureTalk(request):
	context_dict = {}
	return render(request, 'website/architectureTalk.html', context_dict)





def mainSearchCode(request, keyword, tab):
	"""
	  Queries the database for search terms and returns list of results -- used in search and bundles
	"""

	all_results_count = 0
	modules_returned_count = 0
	module_documents_returned_count = 0
	lectures_returned_count = 0
	lecture_segments_returned_count = 0
	lecture_documents_returned_count = 0
	lecture_slides_returned_count = 0
	coming_soon_modules_returned_count = 0
	modules_and_CS_modules_returned_count = 0
	documents_returned_count = 0
	spelling_suggestion = ''
	keyword_compare = keyword.lower()

	if keyword != "":
		kwargs = {
		    "highlight" : {
		        "pre_tags" : ['<span class="highlight">'],
		        "post_tags" : ['</span>'],
		    }
		}
		all_results = SearchQuerySet().auto_query(keyword).highlight(**kwargs)
		spelling_suggestion = all_results.spelling_suggestion()
		modules_returned = []
		module_documents_returned = []
		lectures_returned = []
		lecture_segments_returned = []
		lecture_documents_returned = []
		lecture_slides_returned = []
		coming_soon_modules_returned = []

		# sort search query to bins
		for r in all_results:
			if r.model_name == "modules":
				modules_returned.append(r)
			elif r.model_name == "moduledocuments":
				module_documents_returned.append(r)
			elif r.model_name == "lectures":
				lectures_returned.append(r)
			elif r.model_name == "lecturesegments":
				lecture_segments_returned.append(r)
			elif r.model_name == "lecturedocuments":
				lecture_documents_returned.append(r)
			elif r.model_name == "lectureslides":
				lecture_slides_returned.append(r)
			elif r.model_name == "comingsoonmodules":
				coming_soon_modules_returned.append(r)

		modules_returned_count = len(modules_returned)
		module_documents_returned_count = len(module_documents_returned)
		lectures_returned_count = len(lectures_returned)
		lecture_segments_returned_count = len(lecture_segments_returned)
		lecture_documents_returned_count = len(lecture_documents_returned)
		lecture_slides_returned_count = len(lecture_slides_returned)
		coming_soon_modules_returned_count = len(coming_soon_modules_returned)
		documents_returned_count = module_documents_returned_count + lecture_documents_returned_count + lecture_slides_returned_count
		all_results_count = modules_returned_count + module_documents_returned_count + lectures_returned_count + lecture_segments_returned_count + lecture_documents_returned_count + lecture_slides_returned_count + coming_soon_modules_returned_count

	if all_results_count == 0:
		modules_returned = modules.objects.none()
		moduleDocsCount = moduleDocuments.objects.none()
		lectures_returned = lectures.objects.none()
		lecture_segments_returned = lectureSegments.objects.none()
		lecture_documents_returned = lectureDocuments.objects.none()
		lecture_slides_returned = lectureSlides.objects.none()
		coming_soon_modules_returned = comingSoonModules.objects.none()

		if request.user.is_authenticated():
			# pull bundles
			bundles_returned = bundles.objects.filter(user=request.user)

			# pull user profile
			user_profile = profile.objects.get(user=request.user)

			#attach modules and lectures to profile
			cp_modules = modules.objects.filter(authors_m2m=user_profile).order_by('title')
			user_profile.modules = cp_modules
			cp_lectures = lectures.objects.filter(authors_m2m=user_profile).exclude(extracted=False).order_by('module__title','title')
			user_profile.lectures = cp_lectures
			cp_csmodules = comingSoonModules.objects.filter(authors_m2m=user_profile).order_by('title')
			user_profile.csmodules = cp_csmodules


			# pull saved searches
			saved_searches = savedSearches.objects.filter(user=request.user)

		else:
			bundles_returned = bundles.objects.none()
			user_profile = profile.objects.none()
			saved_searches = savedSearches.objects.none()

		context_dict = {'keyword':keyword, 'modules_returned':modules_returned, 'moduleDocsCount': moduleDocsCount, 'lectures_returned':lectures_returned, 'lecture_segments_returned': lecture_segments_returned, 'lecture_documents_returned':lecture_documents_returned, 'lecture_slides_returned':lecture_slides_returned, 'coming_soon_modules_returned':coming_soon_modules_returned, 'bundles_returned':bundles_returned, 'modules_returned_count':modules_returned_count, 'lectures_returned_count':lectures_returned_count, 'lecture_segments_returned_count':lecture_segments_returned_count, 'lecture_documents_returned_count':lecture_documents_returned_count, 'lecture_slides_returned_count':lecture_slides_returned_count, 'coming_soon_modules_returned_count':coming_soon_modules_returned_count, 'tab': tab, 'user_profile': user_profile, 'saved_searches':saved_searches, 'all_results_count': all_results_count, 'spelling_suggestion': spelling_suggestion, 'keyword_compare':keyword_compare, 'documents_returned_count':documents_returned_count, 'modules_and_CS_modules_returned_count': modules_and_CS_modules_returned_count, 'module_documents_returned':module_documents_returned}

	else:
		# concatonate module querysets
		# first create list of modules
		module_modules = []
		module_documents_modules = []
		lectures_modules = []
		

		for module in modules_returned:
			if module.object is not None:
				module_modules.append(module.object)

		for module_document in module_documents_returned:
			if module_document.object is not None:
				module_document.object.module.highlighted = []
				module_document.object.module.highlighted.append(module_document.highlighted)
				module_documents_modules.append(module_document.object.module)

		for lecture in lectures_returned:
			if lecture.object is not None:
				lecture.object.module.highlighted = []
				lecture.object.module.highlighted.append(lecture.highlighted)
				lectures_modules.append(lecture.object.module)

		module_list = list(chain(module_modules, module_documents_modules, lectures_modules))

		# make unique
		module_set = set(module_list)
		unique_module_list = list(module_set)
		unique_module_list_count = len(unique_module_list)

		#reorder modules if title or author is in the keyword and remove outer list shell from highlighted text
		for module in unique_module_list:

			if module.title.lower().find(keyword.lower()) != -1 or module.authors.lower().find(keyword.lower()) != -1:
				#remove item
				unique_module_list.remove(module)
				#readd item in the front of the list
				unique_module_list.insert(0, module)

			try:
				module.flattened_highlighted = [val for sublist in module.highlighted for val in sublist]
			except AttributeError:
				module.flattened_highlighted = []
			
		#reorder lectures if title or author is in the keyword
		for lec in lectures_returned:
			if lec.object is not None:
				if lec.object.title.lower().find(keyword.lower()) != -1 or lec.object.authors.lower().find(keyword.lower()) != -1:
					#remove item
					lectures_returned.remove(lec)
					#readd item in the front of the list
					lectures_returned.insert(0, lec)

		#reorder lectures if title or author is in the keyword
		for lecdoc in lecture_documents_returned:
			if lecdoc.object is not None:
				if lecdoc.object.title.lower().find(keyword.lower()) != -1 or lecdoc.object.lecture.authors.lower().find(keyword.lower()) != -1:
					#remove item
					lecture_documents_returned.remove(lecdoc)
					#readd item in the front of the list
					lecture_documents_returned.insert(0, lecdoc)

		#reorder lectures if title or author is in the keyword
		for lecslide in lecture_slides_returned:
			if lecslide.object is not None:
				if lecslide.object.lecture.title.lower().find(keyword.lower()) != -1 or lecslide.object.lecture.authors.lower().find(keyword.lower()) != -1:
					#remove item
					lecture_slides_returned.remove(lecslide)
					#readd item in the front of the list
					lecture_slides_returned.insert(0, lecslide)


		#set for template
		modules_returned_unique = unique_module_list
		modules_returned_unique_count = unique_module_list_count

		modules_and_CS_modules_returned_count = modules_returned_unique_count + coming_soon_modules_returned_count

		# sum up all results count again
		all_results_count = modules_returned_unique_count + module_documents_returned_count + lectures_returned_count + lecture_segments_returned_count + lecture_documents_returned_count + lecture_slides_returned_count + coming_soon_modules_returned_count

		if request.user.is_authenticated():
			# pull bundles
			bundles_returned = bundles.objects.filter(user=request.user)

			# pull user profile
			user_profile = profile.objects.get(user=request.user)

			#attach modules and lectures to profiles
			cp_modules = modules.objects.filter(authors_m2m=user_profile).order_by('title')
			user_profile.modules = cp_modules
			cp_lectures = lectures.objects.filter(authors_m2m=user_profile).exclude(extracted=False).order_by('module__title','title')
			user_profile.lectures = cp_lectures
			cp_csmodules = comingSoonModules.objects.filter(authors_m2m=user_profile).order_by('title')
			user_profile.csmodules = cp_csmodules

			# pull saved searches
			saved_searches = savedSearches.objects.filter(user=request.user)

		else:
			bundles_returned = bundles.objects.none()
			user_profile = profile.objects.none()
			saved_searches = savedSearches.objects.none()

		context_dict = {'keyword':keyword, 'modules_returned':modules_returned_unique, 'moduleDocsCount': module_documents_returned_count, 'lectures_returned':lectures_returned, 'lecture_segments_returned':lecture_segments_returned, 'lecture_documents_returned':lecture_documents_returned, 'lecture_slides_returned':lecture_slides_returned, 'coming_soon_modules_returned':coming_soon_modules_returned, 'bundles_returned':bundles_returned, 'modules_returned_count':modules_returned_unique_count, 'lectures_returned_count':lectures_returned_count, 'lecture_segments_returned_count':lecture_segments_returned_count ,'lecture_documents_returned_count':lecture_documents_returned_count, 'lecture_slides_returned_count':lecture_slides_returned_count, 'coming_soon_modules_returned_count':coming_soon_modules_returned_count, 'tab': tab, 'user_profile': user_profile, 'saved_searches':saved_searches, 'all_results_count': all_results_count, 'spelling_suggestion': spelling_suggestion, 'keyword_compare':keyword_compare, 'documents_returned_count':documents_returned_count, 'modules_and_CS_modules_returned_count': modules_and_CS_modules_returned_count, 'module_documents_returned':module_documents_returned}

	return context_dict


def search(request):


	"""
	  Queries the database for search terms and returns list of results
	"""
	if request.method == 'POST':
		# search terms
		try:
			keyword = request.POST['keyword']
			tab = 'search'
			context_dict = mainSearchCode(request, keyword, tab)
		except KeyError:
			return HttpResponseRedirect("/")

	else:
		try:
			keyword = request.GET['keyword']
			tab = 'search'
			context_dict = mainSearchCode(request, keyword, tab)
		except KeyError:
			return HttpResponseRedirect("/")


	return render(request, 'website/search.html', context_dict)


@login_required
def mybundles(request):

	# which url route was this from
	if request.path == '/profile/':
		tab = 'profile'
	elif request.path == '/bundles/':
		tab = 'bundle'
	else:
		tab = 'searches'

	"""
	  Pulls all user bundles and prints
	"""
	if request.user.is_authenticated():
		# pull bundles
		bundles_returned = bundles.objects.filter(user=request.user)

		for bundle_returned in bundles_returned:
			#start collapsed
			bundle_returned.in_class = ''
			bundle_returned.show_more = 'SHOW MORE'

			#get return all other content
			bundle_returned.bundle_modules = bundleModule.objects.filter(bundle=bundle_returned)
			bundle_returned.bundle_lectures = bundleLecture.objects.filter(bundle=bundle_returned)
			bundle_returned.bundle_lecture_documents = bundleLectureDocument.objects.filter(bundle=bundle_returned)
			bundle_returned.bundle_lecture_slides = bundleLectureSlides.objects.filter(bundle=bundle_returned)

			# for each module in a bundle, attach the module documents and the lectures
			for bundle in bundle_returned.bundle_modules:
				#look up module docs
				moduleDocs = moduleDocuments.objects.filter(module=bundle.module)
				# just get the file name
				for doc in moduleDocs:
					document = str(doc.document)
					document = document.split('/')
					doc.documentName = document[2]

				#attach to the bundle
				bundle.module.moduleDocs = moduleDocs

				#look up the lectures
				moduleLecs = lectures.objects.filter(module=bundle.module).exclude(extracted=False).order_by('title')
				# get the file name
				for lec in moduleLecs:

					# order by lecture number
					first_word = lec.title.strip().lower().split(' ', 1)[0]
					if first_word == 'lecture':
						first_number = lec.title.strip().lower().split(' ')[1].replace('.','').replace(':','')
						try:
							lec.numeric_order = int(first_number)
						except ValueError:
							lec.numeric_order = 0
					else:
						lec.numeric_order = 0

					lecture = str(lec.presentation)
					lecture = lecture.split('/')
					lec.lectureName = lecture[2]
					# get lecture documents
					lectureDocs = lectureDocuments.objects.filter(lecture=lec)
					lec.lectureDocs = lectureDocs
					for lecDoc in lec.lectureDocs:
						document = str(lecDoc.document)
						document = document.split('/')
						lecDoc.documentName = document[2]

				# order titles minus articles
				moduleLecs_ordered = sorted(moduleLecs, key=operator.attrgetter('numeric_order'))

				bundle.module.lectures = moduleLecs_ordered

			#for each lecture, attach lecture docs
			for bundle in bundle_returned.bundle_lectures:
				# get lecture documents
				lectureDocs = lectureDocuments.objects.filter(lecture=bundle.lecture)
				bundle.lecture.lectureDocs = lectureDocs
				for lecDoc in bundle.lecture.lectureDocs:
					document = str(lecDoc.document)
					document = document.split('/')
					lecDoc.documentName = document[2]


			#for each lecture document strip out name of file
			for bundle in bundle_returned.bundle_lecture_documents:
				document = str(bundle.lectureDocument.document)
				document = document.split('/')
				bundle.lectureDocument.documentName = document[2]

			#for each lecture slide strip out name of file and slide notes file
			for bundle in bundle_returned.bundle_lecture_slides:
				slide = str(bundle.lectureSlide.presentation)
				slide = slide.split('/')
				bundle.lectureSlide.slideName = slide[2]


	else:
		bundles_returned = bundles.objects.none()


	# also pull saved searches
	if request.user.is_authenticated():
		# pull saved searches
		saved_searches = savedSearches.objects.filter(user=request.user)
	else:
		saved_searches = savedSearches.objects.none()

	# also pull my profile
	if request.user.is_authenticated():
		# pull user profile
		user_profile = profile.objects.get(user=request.user)

		#attach modules and lectures to profile
		cp_modules = modules.objects.filter(authors_m2m=user_profile).order_by('title')
		user_profile.modules = cp_modules
		cp_lectures = lectures.objects.filter(authors_m2m=user_profile).exclude(extracted=False).order_by('module__title','title')
		user_profile.lectures = cp_lectures
		cp_csmodules = comingSoonModules.objects.filter(authors_m2m=user_profile).order_by('title')
		user_profile.csmodules = cp_csmodules

	else:
		user_profile = profile.objects.none()

	context_dict = {'bundles_returned':bundles_returned, 'saved_searches':saved_searches, 'user_profile': user_profile, 'tab':tab}


	return render(request, 'website/bundle_list.html', context_dict)


# @login_required
# def myprofile(request):


# 	"""
# 	  Queries the database for search terms and returns list of results; goes to bundle
# 	"""

# 	# placeholder keyword to return all items
# 	tab = 'profile'

# 	context_dict = {}

# 	return render(request, 'website/profile.html', context_dict)


# @login_required
# def mysavedsearches(request):

# 	"""
# 	  Pulls all user searches and prints
# 	"""
# 	if request.user.is_authenticated():

# 		# pull saved searches
# 		saved_searches = savedSearches.objects.filter(user=request.user)


# 	else:
# 		saved_searches = savedSearches.objects.none()

# 	context_dict = {'saved_searches':saved_searches,}


# 	return render(request, 'website/saved_searches.html', context_dict)



def showModule(request, id=None):
	"""
	  Response from AJAX request to show module in sidebar
	"""
	#get module
	module_returned = modules.objects.get(pk=id)

	#look up module docs
	moduleDocs = moduleDocuments.objects.filter(module=module_returned)
	# just get the file name
	for doc in moduleDocs:
		document = str(doc.document)
		document = document.split('/')
		doc.documentName = document[2]

	#look up the lectures
	moduleLecs = lectures.objects.filter(module=module_returned).exclude(extracted=False).order_by('title')
	# get the file name
	for lec in moduleLecs:

		# order by lecture number
		first_word = lec.title.strip().lower().split(' ', 1)[0]
		if first_word == 'lecture':
			first_number = lec.title.strip().lower().split(' ')[1].replace('.','').replace(':','')
			try:
				lec.numeric_order = int(first_number)
			except ValueError:
				lec.numeric_order = 0
		else:
			lec.numeric_order = 0

		lecture = str(lec.presentation)
		lecture = lecture.split('/')
		lec.lectureName = lecture[2]
		# get lecture documents
		lectureDocs = lectureDocuments.objects.filter(lecture=lec)
		lec.lectureDocs = lectureDocs
		for lecDoc in lec.lectureDocs:
			document = str(lecDoc.document)
			document = document.split('/')
			lecDoc.documentName = document[2]

	# order titles minus articles
	moduleLecs_ordered = sorted(moduleLecs, key=operator.attrgetter('numeric_order'))

	moduleDocsCount = moduleDocuments.objects.filter(module=module_returned)
	contents = []
	for doc in moduleDocsCount:
		#add document contents to modules
		contents.append(doc.document_contents)

	# join document contents together
	module_returned.document_contents = '\n'.join(contents)

	#get first lecture uploaded
	earliest_lecture = lectures.objects.filter(module=module_returned).exclude(extracted=False).earliest('created')

	if request.user.is_authenticated():
		# pull bundles
		bundles_returned = bundles.objects.filter(user=request.user)
	else:
		bundles_returned = bundles.objects.none()

	# get comments
	comments = modulesComments.objects.filter(module=module_returned)

	# comment form
	comment_form = modulesCommentsForm()

	context_dict = {'module_returned':module_returned, 'earliest_lecture':earliest_lecture, 'bundles_returned':bundles_returned, 'moduleDocs':moduleDocs, 'moduleLecs':moduleLecs_ordered, 'comments':comments , 'comment_form':comment_form}
	return render(request, 'website/show_module.html', context_dict)

def showLecture(request, id=None):
	"""
	  Response from AJAX request to show lecture in sidebar
	"""
	#get lecture
	lecture_returned = lectures.objects.get(pk=id)
	lecture_slides = lectureSlides.objects.filter(lecture=lecture_returned).order_by("slide_number")

	if request.user.is_authenticated():
		# pull bundles
		bundles_returned = bundles.objects.filter(user=request.user)
	else:
		bundles_returned = bundles.objects.none()

	# get comments
	comments = lecturesComments.objects.filter(lecture=lecture_returned)

	# comment form
	comment_form = lecturesCommentsForm()

	context_dict = {'lecture_returned':lecture_returned, 'lecture_slides': lecture_slides, 'bundles_returned':bundles_returned, 'comments':comments , 'comment_form':comment_form}
	return render(request, 'website/show_lecture.html', context_dict)

def showLectureSegment(request, id=None):
	"""
	  Response from AJAX request to show lecture segment in sidebar
	"""
	#get lecture
	lecture_segment_returned = lectureSegments.objects.get(pk=id)
	lecture_slides = lectureSlidesSegment.objects.filter(lecture_segment=lecture_segment_returned).order_by("slide_number")

	if request.user.is_authenticated():
		# pull bundles
		bundles_returned = bundles.objects.filter(user=request.user)
	else:
		bundles_returned = bundles.objects.none()

	# get comments
	comments = lectureSegmentsComments.objects.filter(lectureSegment=lecture_segment_returned)

	# comment form
	comment_form = lectureSegmentsCommentsForm()

	context_dict = {'lecture_segment_returned':lecture_segment_returned, 'lecture_slides': lecture_slides, 'bundles_returned':bundles_returned, 'comments':comments , 'comment_form':comment_form}
	return render(request, 'website/show_lecture_segment.html', context_dict)

def showLectureDocument(request, id=None):
	"""
	  Response from AJAX request to show lecture document in sidebar
	"""
	#get lecture document
	lecture_document_returned = lectureDocuments.objects.get(pk=id)

	if request.user.is_authenticated():
		# pull bundles
		bundles_returned = bundles.objects.filter(user=request.user)
	else:
		bundles_returned = bundles.objects.none()

	# get comments
	comments = lectureDocumentsComments.objects.filter(lectureDocument=lecture_document_returned)

	# comment form
	comment_form = lectureDocumentsCommentsForm()

	context_dict = {'lecture_document_returned':lecture_document_returned, 'bundles_returned':bundles_returned, 'comments':comments , 'comment_form':comment_form}
	return render(request, 'website/show_lecture_document.html', context_dict)

def showLectureSlide(request, id=None):
	"""
	  Response from AJAX request to show lecture slide in sidebar
	"""
	#get lecture slide
	lecture_slide_returned = lectureSlides.objects.get(pk=id)

	if request.user.is_authenticated():
		# pull bundles
		bundles_returned = bundles.objects.filter(user=request.user)
	else:
		bundles_returned = bundles.objects.none()

	# get comments
	comments = lectureSlidesComments.objects.filter(lectureSlide=lecture_slide_returned)

	# comment form
	comment_form = lectureSlidesCommentsForm()

	context_dict = {'lecture_slide_returned':lecture_slide_returned, 'bundles_returned':bundles_returned, 'comments':comments , 'comment_form':comment_form}
	return render(request, 'website/show_lecture_slide.html', context_dict)

def showLectureModal(request, id=None):
	"""
	  Response from AJAX request to show lecture slides in modal
	"""
	#get lecture
	lecture_returned = lectures.objects.get(pk=id)
	lecture_slides = lectureSlides.objects.filter(lecture=lecture_returned).order_by("slide_number")

	context_dict = {'lecture_returned':lecture_returned, 'lecture_slides':lecture_slides}
	return render(request, 'website/show_lecture_modal.html', context_dict)

def showModuleDescription(request, id=None):
  """
    Response from AJAX request to show module description in modal
  """
  #get lecture
  module_returned = modules.objects.get(pk=id)

  context_dict = {'module_returned':module_returned}
  return render(request, 'website/show_module_description.html', context_dict)

def showLectureDescription(request, id=None):
  """
    Response from AJAX request to show lecture description in modal
  """
  #get lecture
  lecture_returned = lectures.objects.get(pk=id)

  context_dict = {'lecture_returned':lecture_returned}
  return render(request, 'website/show_lecture_description.html', context_dict)

def showMemberIntroduction(request, id=None):
  """
    Response from AJAX request to show member introduction in modal
  """
  #get profile
  profile_returned = profile.objects.get(pk=id)

  context_dict = {'profile_returned':profile_returned}
  return render(request, 'website/show_member_introduction.html', context_dict)


def showLectureSegmentModal(request, id=None):
	"""
	  Response from AJAX request to show lecture segment slides in modal
	"""
	#get lecture
	lecture_returned = lectureSegments.objects.get(pk=id)
	lecture_slides = lectureSlidesSegment.objects.filter(lecture_segment=lecture_returned).order_by("slide_number")

	context_dict = {'lecture_returned':lecture_returned, 'lecture_slides':lecture_slides}
	return render(request, 'website/show_lecture_modal.html', context_dict)

def showLectureSlideModal(request, id=None):
	"""
	  Response from AJAX request to show lecture slide in modal
	"""
	#get lecture
	lecture_slide = lectureSlides.objects.get(pk=id)

	context_dict = {'lecture_slide':lecture_slide}
	return render(request, 'website/show_lecture_slide_modal.html', context_dict)

def createNewBundle(request):
	"""
	  AJAX request to create new bundle and attach an item to the bundle
	"""

	if request.method == 'GET':
		#gather variables from get request
		title = request.GET.get("title","")
		itemid = request.GET.get("itemid","")

		# create bundle
		bundle = bundles(user=request.user, title=title)
		bundle.save()

		# add item to appropriate bundle
		itemid = itemid.split("_")
		if itemid[0] == 'module':
			#look up module
			module = modules.objects.get(pk=itemid[1])
			bm = bundleModule(bundle=bundle, module=module)
			bm.save()
		elif itemid[0] == 'lecture':
			#look up lecture
			lecture = lectures.objects.get(pk=itemid[1])
			bl = bundleLecture(bundle=bundle, lecture=lecture)
			bl.save()
		elif itemid[0] == 'lecturesegment':
			#look up lecture segment
			lectureSegment = lectureSegments.objects.get(pk=itemid[1])
			blseg = bundleLectureSegments(bundle=bundle, lectureSegment=lectureSegment)
			blseg.save()
		elif itemid[0] == 'lecturedocument':
			#look up lectureDocument
			lectureDocument = lectureDocuments.objects.get(pk=itemid[1])
			bld = bundleLectureDocument(bundle=bundle, lectureDocument=lectureDocument)
			bld.save()
		else:
			#look up lectureSlide
			lectureSlide = lectureSlides.objects.get(pk=itemid[1])
			bls = bundleLectureSlides(bundle=bundle, lectureSlide=lectureSlide)
			bls.save()

	# pull bundles
	bundles_returned = bundles.objects.filter(user=request.user)


	context_dict = {'bundles_returned':bundles_returned}
	return render(request, 'website/bundle_dropdown.html', context_dict)


def addToBundle(request):
	"""
	  AJAX request to attach an item to the bundle
	"""

	if request.method == 'GET':
		#gather variables from get request
		bundleid = request.GET.get("bundle","")
		itemid = request.GET.get("itemid","")

		# get bundle
		bundle = bundles.objects.get(pk=bundleid)

		# add item to appropriate bundle
		itemid = itemid.split("_")
		if itemid[0] == 'module':
			#look up module
			module = modules.objects.get(pk=itemid[1])
			#does bundle/module exist already?
			if bundleModule.objects.filter(bundle=bundle, module=module).exists() == False:
				bm = bundleModule(bundle=bundle, module=module)
				bm.save()
		elif itemid[0] == 'lecture':
			#look up lecture
			lecture = lectures.objects.get(pk=itemid[1])
			if bundleLecture.objects.filter(bundle=bundle, lecture=lecture).exists() == False:
				bl = bundleLecture(bundle=bundle, lecture=lecture)
				bl.save()
		elif itemid[0] == 'lecturesegment':
			#look up lectureDocument
			lectureSegment = lectureSegments.objects.get(pk=itemid[1])
			if bundleLectureSegments.objects.filter(bundle=bundle, lectureSegment=lectureSegment).exists() == False:
				blseg = bundleLectureSegments(bundle=bundle, lectureSegment=lectureSegment)
				blseg.save()
		elif itemid[0] == 'lecturedocument':
			#look up lectureDocument
			lectureDocument = lectureDocuments.objects.get(pk=itemid[1])
			if bundleLectureDocument.objects.filter(bundle=bundle, lectureDocument=lectureDocument).exists() == False:
				bld = bundleLectureDocument(bundle=bundle, lectureDocument=lectureDocument)
				bld.save()
		else:
			#look up lectureSlide
			lectureSlide = lectureSlides.objects.get(pk=itemid[1])
			if bundleLectureSlides.objects.filter(bundle=bundle, lectureSlide=lectureSlide).exists() == False:
				bls = bundleLectureSlides(bundle=bundle, lectureSlide=lectureSlide)
				bls.save()

	# pull bundles
	bundles_returned = bundles.objects.filter(user=request.user)


	context_dict = {'bundles_returned':bundles_returned}
	return render(request, 'website/bundle_dropdown.html', context_dict)


def removeFromBundle(request):
	"""
	  AJAX request to remove an item to the bundle
	"""

	if request.method == 'GET':
		#gather variables from get request
		bundleid = request.GET.get("bundle","")
		itemid = request.GET.get("itemid","")
		typename = request.GET.get("type","")

		# get bundle
		bundle = bundles.objects.get(pk=bundleid)

		# remove item from appropriate bundle
		if typename == 'module':
			#look up module
			module = modules.objects.get(pk=itemid)
			bm = bundleModule.objects.filter(bundle=bundle, module=module)
			for rec in bm:
				rec.delete()

		elif typename == 'lecture':
			#look up lecture
			lecture = lectures.objects.get(pk=itemid)
			bl = bundleLecture.objects.filter(bundle=bundle, lecture=lecture)
			for rec in bl:
				rec.delete()

		elif typename == 'lecturesegment':
			#look up lectureDocument
			lectureSegment = lectureSegments.objects.get(pk=itemid)
			blseg = bundleLectureSegments.objects.filter(bundle=bundle, lectureSegment=lectureSegment)
			for rec in blseg:
				rec.delete()

		elif typename == 'lecturedocument':
			#look up lectureDocument
			lectureDocument = lectureDocuments.objects.get(pk=itemid)
			bld = bundleLectureDocument.objects.filter(bundle=bundle, lectureDocument=lectureDocument)
			for rec in bld:
				rec.delete()
		else:
			#look up lectureSlide
			lectureSlide = lectureSlides.objects.get(pk=itemid)
			bls = bundleLectureSlides.objects.filter(bundle=bundle, lectureSlide=lectureSlide)
			for rec in bls:
				rec.delete()

	# pull bundles
	bundles_returned = bundles.objects.filter(user=request.user)


	context_dict = {'bundles_returned':bundles_returned}
	return render(request, 'website/bundle_dropdown.html', context_dict)

def removeBundle(request):
	"""
	  AJAX request to remove a bundle
	"""

	if request.method == 'GET':
		#gather variables from get request
		bundleid = request.GET.get("bundle","")

		# get bundle and remove it
		bundle = bundles.objects.get(pk=bundleid).delete()


	# pull remaining bundles
	bundles_returned = bundles.objects.filter(user=request.user)

	context_dict = {'bundles_returned':bundles_returned}
	return render(request, 'website/bundle_dropdown.html', context_dict)


def removeSearch(request):

	"""
	  AJAX request to save search string
	"""

	if request.method == 'GET':
		#gather variables from get request
		searchid = request.GET.get("search","")

		# create saved search unless one already exists
		search = savedSearches.objects.get(user=request.user, pk=searchid).delete()

	saved_searches = savedSearches.objects.filter(user=request.user)

	context_dict = {'saved_searches':saved_searches}
	return render(request, 'website/current_saved_searches.html', context_dict)


def showBundle(request, id=None):
	"""
	  Response from AJAX request to show bundle in sidebar
	"""
	#get lecture slide
	bundle_returned = bundles.objects.get(pk=id, user=request.user)

	#get return all other content
	bundle_returned.bundle_modules = bundleModule.objects.filter(bundle=bundle_returned)
	bundle_returned.bundle_lectures = bundleLecture.objects.filter(bundle=bundle_returned)
	bundle_returned.bundle_lecture_segments = bundleLectureSegments.objects.filter(bundle=bundle_returned)
	bundle_returned.bundle_lecture_documents = bundleLectureDocument.objects.filter(bundle=bundle_returned)
	bundle_returned.bundle_lecture_slides = bundleLectureSlides.objects.filter(bundle=bundle_returned)

	# for each module in a bundle, attach the module documents and the lectures
	for bundle in bundle_returned.bundle_modules:
		#look up module docs
		moduleDocs = moduleDocuments.objects.filter(module=bundle.module)
		# just get the file name
		for doc in moduleDocs:
			document = str(doc.document)
			document = document.split('/')
			doc.documentName = document[2]

		#attach to the bundle
		bundle.module.moduleDocs = moduleDocs

		#look up the lectures
		moduleLecs = lectures.objects.filter(module=bundle.module).exclude(extracted=False).order_by('title')
		# get the file name
		for lec in moduleLecs:

			# order by lecture number
			first_word = lec.title.strip().lower().split(' ', 1)[0]
			if first_word == 'lecture':
				first_number = lec.title.strip().lower().split(' ')[1].replace('.','').replace(':','')
				try:
					lec.numeric_order = int(first_number)
				except ValueError:
					lec.numeric_order = 0
			else:
				lec.numeric_order = 0

			lecture = str(lec.presentation)
			lecture = lecture.split('/')
			lec.lectureName = lecture[2]
			# get lecture documents
			lectureDocs = lectureDocuments.objects.filter(lecture=lec)
			lec.lectureDocs = lectureDocs
			for lecDoc in lec.lectureDocs:
				document = str(lecDoc.document)
				document = document.split('/')
				lecDoc.documentName = document[2]

		# order titles minus articles
		moduleLecs_ordered = sorted(moduleLecs, key=operator.attrgetter('numeric_order'))

		bundle.module.lectures = moduleLecs_ordered


	#for each lecture, attach lecture docs
	for bundle in bundle_returned.bundle_lectures:
		# get lecture documents
		lectureDocs = lectureDocuments.objects.filter(lecture=bundle.lecture)
		bundle.lecture.lectureDocs = lectureDocs
		for lecDoc in bundle.lecture.lectureDocs:
			document = str(lecDoc.document)
			document = document.split('/')
			lecDoc.documentName = document[2]
		
	#for each lecture document strip out name of file
	for bundle in bundle_returned.bundle_lecture_documents:
		document = str(bundle.lectureDocument.document)
		document = document.split('/')
		bundle.lectureDocument.documentName = document[2]

	#for each lecture slide strip out name of file and slide notes file
	for bundle in bundle_returned.bundle_lecture_slides:
		slide = str(bundle.lectureSlide.presentation)
		slide = slide.split('/')
		bundle.lectureSlide.slideName = slide[2]

	bundle_returned.in_class = 'in'
	bundle_returned.show_more = 'SHOW LESS'


	context_dict = {'bundle_returned':bundle_returned,}
	return render(request, 'website/show_bundle.html', context_dict)


def zipUpBundle(request, id=None):
	"""
	  Response from AJAX request to zip up bundle and create download
	"""
	#folder for zip file
	folder = "/zip_files/"+ id +"/"
	filename = "GAHTC_bundle_"+ id +".zip"
	zipfolder = "GAHTC_bundle_"+ id

	if not os.path.exists(MEDIA_ROOT + folder):
		os.makedirs(MEDIA_ROOT + folder)

	#create zip file
	with zipfile.ZipFile(MEDIA_ROOT + folder + filename, "w", allowZip64=True) as myzip:

		#get bundle
		bundle_returned = bundles.objects.get(pk=id, user=request.user)

		# set downloaded to be true
		bundle_returned.downloaded = True
		bundle_returned.save()

		#get return all other content
		bundle_modules = bundleModule.objects.filter(bundle=bundle_returned)
		bundle_lectures = bundleLecture.objects.filter(bundle=bundle_returned)
		bundle_lecture_segments = bundleLectureSegments.objects.filter(bundle=bundle_returned)
		bundle_lecture_documents = bundleLectureDocument.objects.filter(bundle=bundle_returned)
		bundle_lecture_slides = bundleLectureSlides.objects.filter(bundle=bundle_returned)

		# for each module in a bundle, attach the module documents and the lectures
		for bundle in bundle_modules:
			# create directory for module

			#look up module docs
			moduleDocs = moduleDocuments.objects.filter(module=bundle.module)
			#loop over docs and add to zip archive in the correct folder
			for doc in moduleDocs:
				modTitle = bundle.module.title.replace("<", "").replace(">", "").replace(":", "").replace('"', '').replace("/", "").replace("\\", "").replace("|", "").replace("?", "").replace("*", "").replace(",", "").replace(".", "")
				"".join(c for c in modTitle if c.isalnum() or c==' ').rstrip()
				modTitle = '_'.join(modTitle.split())[:30]

				document = unicode(doc.document)
				document = document.split('/')
				doc_name = document[2].encode('utf8', 'replace')
				directory = os.path.join(zipfolder+ "/Modules/" + modTitle + "/Documents", doc_name.decode('utf8', 'replace'))
				myzip.write(MEDIA_ROOT + '/' + str(doc.document), directory)

			#look up the lectures
			moduleLecs = lectures.objects.filter(module=bundle.module).exclude(extracted=False).order_by('title')
			#loop over lectures and add to zip archive in the correct folder
			for i, lec in enumerate(moduleLecs,1):
				modTitle = bundle.module.title.replace("<", "").replace(">", "").replace(":", "").replace('"', '').replace("/", "").replace("\\", "").replace("|", "").replace("?", "").replace("*", "").replace(",", "").replace(".", "")
				"".join(c for c in modTitle if c.isalnum() or c==' ').rstrip()
				modTitle = '_'.join(modTitle.split())[:30]

				lecTitle = lec.title.replace("<", "").replace(">", "").replace(":", "").replace('"', '').replace("/", "").replace("\\", "").replace("|", "").replace("?", "").replace("*", "").replace(",", "").replace(".", "")
				"".join(c for c in lecTitle if c.isalnum() or c==' ').rstrip()
				lecTitle = '_'.join(lecTitle.split())[:30]

				document = unicode(lec.presentation)
				document = document.split('/')
				doc_name = document[2].encode('utf8', 'replace')
				directory = os.path.join(zipfolder+ "/Modules/" + modTitle + "/Lectures/" + lecTitle, doc_name.decode('utf8', 'replace'))
				myzip.write(MEDIA_ROOT + '/' + str(lec.presentation), directory)

				# look up lecture documents
				moduleLecDocs = lectureDocuments.objects.filter(lecture=lec)
				for lecDoc in moduleLecDocs:
					document = unicode(lecDoc.document)
					document = document.split('/')
					doc_name = document[2].encode('utf8', 'replace')
					directory = os.path.join(zipfolder+ "/Modules/" + modTitle + "/Lectures/" + lecTitle, doc_name.decode('utf8', 'replace'))
					myzip.write(MEDIA_ROOT + '/' + str(lecDoc.document), directory)


		# for each lecture write to zip archive
		for i, bundle in enumerate(bundle_lectures,1):
			lecTitle = bundle.lecture.title.replace("<", "").replace(">", "").replace(":", "").replace('"', '').replace("/", "").replace("\\", "").replace("|", "").replace("?", "").replace("*", "").replace(",", "").replace(".", "")
			"".join(c for c in lecTitle if c.isalnum() or c==' ').rstrip()
			lecTitle = '_'.join(lecTitle.split())[:30]			

			document = unicode(bundle.lecture.presentation)
			document = document.split('/')
			doc_name = document[2].encode('utf8', 'replace')
			directory = os.path.join(zipfolder+ "/Lectures/" + lecTitle, doc_name.decode('utf8', 'replace'))
			myzip.write(MEDIA_ROOT + '/' + str(bundle.lecture.presentation), directory)

			# look up lecture documents and add these to zip archive
			lecDocs = lectureDocuments.objects.filter(lecture=bundle.lecture)
			for lecDoc in lecDocs:
				document = unicode(lecDoc.document)
				document = document.split('/')
				doc_name = document[2].encode('utf8', 'replace')
				directory = os.path.join(zipfolder+ "/Lectures/" + lecTitle, doc_name.decode('utf8', 'replace'))
				myzip.write(MEDIA_ROOT + '/' + str(lecDoc.document), directory)


		# for each lecture segment write to zip archive
		for i, bundle in enumerate(bundle_lecture_segments,1):
			lecTitle = bundle.lectureSegment.title.replace("<", "").replace(">", "").replace(":", "").replace('"', '').replace("/", "").replace("\\", "").replace("|", "").replace("?", "").replace("*", "").replace(",", "").replace(".", "")
			"".join(c for c in lecTitle if c.isalnum() or c==' ').rstrip()
			lecTitle = '_'.join(lecTitle.split())[:30]	

			document = unicode(bundle.lectureSegment.presentation)
			document = document.split('/')
			doc_name = document[2].encode('utf8', 'replace')
			directory = os.path.join(zipfolder+ "/Lecture_segments/" + lecTitle, doc_name.decode('utf8', 'replace'))
			myzip.write(MEDIA_ROOT + '/' + str(bundle.lectureSegment.presentation), directory)


		#for each lecture document strip out name of file
		for i, bundle in enumerate(bundle_lecture_documents,1):	
			lecTitle = bundle.lectureDocument.title.replace("<", "").replace(">", "").replace(":", "").replace('"', '').replace("/", "").replace("\\", "").replace("|", "").replace("?", "").replace("*", "").replace(",", "").replace(".", "")
			"".join(c for c in lecTitle if c.isalnum() or c==' ').rstrip()
			lecTitle = '_'.join(lecTitle.split())[:30]	

			document = unicode(bundle.lectureDocument.document)
			document = document.split('/')
			doc_name = document[2].encode('utf8', 'replace')
			directory = os.path.join(zipfolder+ "/Lecture_documents/" + lecTitle, doc_name.decode('utf8', 'replace'))
			myzip.write(MEDIA_ROOT + '/' + str(bundle.lectureDocument.document), directory)

		#for each lecture slide strip out name of file and slide notes file
		for i, bundle in enumerate(bundle_lecture_slides,1):
			lecTitle = bundle.lectureSlide.lecture.title.replace("<", "").replace(">", "").replace(":", "").replace('"', '').replace("/", "").replace("\\", "").replace("|", "").replace("?", "").replace("*", "").replace(",", "").replace(".", "")
			"".join(c for c in lecTitle if c.isalnum() or c==' ').rstrip()
			lecTitle = '_'.join(lecTitle.split())[:30]	

			document = unicode(bundle.lectureSlide.presentation)
			document = document.split('/')
			doc_name = document[2].encode('utf8', 'replace')
			directory = os.path.join(zipfolder+ "/Lecture_slides/" + lecTitle, doc_name.decode('utf8', 'replace'))
			myzip.write(MEDIA_ROOT + '/' + str(bundle.lectureSlide.presentation), directory)

	#get size of zip file
	filesize = os.path.getsize(MEDIA_ROOT + folder + filename)

	context_dict = {'folder':folder, 'filename':filename, 'filesize':filesize, 'bundleid':id}
	return render(request, 'website/bundle_download.html', context_dict)


def zipUpModule(request, id=None):
	"""
	  Response from AJAX request to zip up module and create download
	"""
	# get module
	module_returned = modules.objects.get(pk=id)
	modTitle = module_returned.title.replace("<", "").replace(">", "").replace(":", "").replace('"', '').replace("/", "").replace("\\", "").replace("|", "").replace("?", "").replace("*", "").replace(",", "").replace(".", "")
	"".join(c for c in modTitle if c.isalnum() or c==' ').rstrip()
	modTitle = '_'.join(modTitle.split())[:30]

	#folder for zip file
	folder = "/zip_files/module_"+ id +"/"
	filename = "GAHTC_"+ modTitle +".zip"
	zipfolder = "GAHTC_"+ modTitle

	if not os.path.exists(MEDIA_ROOT + folder):
		os.makedirs(MEDIA_ROOT + folder)

	#create zip file
	with zipfile.ZipFile(MEDIA_ROOT + folder + filename, "w", allowZip64=True) as myzip:

		# mark as downloaded
		d = userModuleDownload()
		d.user = request.user
		d.module = module_returned
		d.downloaded = True
		d.save()

		#look up module docs
		moduleDocs = moduleDocuments.objects.filter(module=module_returned)
		#loop over docs and add to zip archive in the correct folder
		for doc in moduleDocs:
			document = unicode(doc.document)
			document = document.split('/')
			doc_name = document[2].encode('utf8', 'replace')
			directory = os.path.join(zipfolder+ "/Documents", doc_name.decode('utf8', 'replace'))
			myzip.write(MEDIA_ROOT + '/' + str(doc.document), directory)

		#look up the lectures
		moduleLecs = lectures.objects.filter(module=module_returned).exclude(extracted=False).order_by('title')
		#loop over lectures and add to zip archive in the correct folder
		for i, lec in enumerate(moduleLecs,1):
			lecTitle = lec.title.replace("<", "").replace(">", "").replace(":", "").replace('"', '').replace("/", "").replace("\\", "").replace("|", "").replace("?", "").replace("*", "").replace(",", "").replace(".", "")
			"".join(c for c in lecTitle if c.isalnum() or c==' ').rstrip()
			lecTitle = '_'.join(lecTitle.split())[:30]
			document = unicode(lec.presentation)
			document = document.split('/')
			doc_name = document[2].encode('utf8', 'replace')
			directory = os.path.join(zipfolder+ "/Lectures/" + lecTitle, doc_name.decode('utf8', 'replace'))
			myzip.write(MEDIA_ROOT + '/' + str(lec.presentation), directory)

			# look up lecture documents
			moduleLecDocs = lectureDocuments.objects.filter(lecture=lec)
			for lecDoc in moduleLecDocs:
				document = unicode(lecDoc.document)
				document = document.split('/')
				doc_name = document[2].encode('utf8', 'replace')
				directory = os.path.join(zipfolder+ "/Lectures/" + lecTitle, doc_name.decode('utf8', 'replace'))
				myzip.write(MEDIA_ROOT + '/' + str(lecDoc.document), directory)



	#get size of zip file
	filesize = os.path.getsize(MEDIA_ROOT + folder + filename)

	context_dict = {'folder':folder, 'filename':filename, 'filesize':filesize, 'moduleid':id}
	return render(request, 'website/module_download.html', context_dict)


def zipUpLecture(request, id=None):
	"""
	  Response from AJAX request to zip up lectures and create download
	"""

	lec = lectures.objects.get(pk=id)
	lecTitle = lec.title.replace("<", "").replace(">", "").replace(":", "").replace('"', '').replace("/", "").replace("\\", "").replace("|", "").replace("?", "").replace("*", "").replace(",", "").replace(".", "")
	"".join(c for c in lecTitle if c.isalnum() or c==' ').rstrip()
	lecTitle = '_'.join(lecTitle.split())[:30]	

	#folder for zip file
	folder = "/zip_files/lecture_"+ id +"/"
	filename = "GAHTC_"+ lecTitle +".zip"
	zipfolder = "GAHTC_"+ lecTitle

	if not os.path.exists(MEDIA_ROOT + folder):
		os.makedirs(MEDIA_ROOT + folder)

	#create zip file
	with zipfile.ZipFile(MEDIA_ROOT + folder + filename, "w", allowZip64=True) as myzip:

		# mark as downloaded
		d = userLectureDownload()
		d.user = request.user
		d.lecture = lec
		d.downloaded = True
		d.save()

		document = unicode(lec.presentation)
		document = document.split('/')
		doc_name = document[2].encode('utf8', 'replace')
		directory = os.path.join(zipfolder, doc_name.decode('utf8', 'replace'))
		myzip.write(MEDIA_ROOT + '/' + str(lec.presentation), directory)

		# look up lecture documents
		lecDocs = lectureDocuments.objects.filter(lecture=lec)
		for lecDoc in lecDocs:
			document = unicode(lecDoc.document)
			document = document.split('/')
			doc_name = document[2].encode('utf8', 'replace')
			directory = os.path.join(zipfolder, doc_name.decode('utf8', 'replace'))
			myzip.write(MEDIA_ROOT + '/' + str(lecDoc.document), directory)



	#get size of zip file
	filesize = os.path.getsize(MEDIA_ROOT + folder + filename)

	context_dict = {'folder':folder, 'filename':filename, 'filesize':filesize, 'lectureid':id}
	return render(request, 'website/lecture_download.html', context_dict)



def refreshSidebarBundle(request):

	# pull bundles
	bundles_returned = bundles.objects.filter(user=request.user)

	for bundle_returned in bundles_returned:
		#start collapesed
		bundle_returned.in_class = ''
		bundle_returned.show_more = 'SHOW MORE'

		#get return all other content
		bundle_returned.bundle_modules = bundleModule.objects.filter(bundle=bundle_returned)
		bundle_returned.bundle_lectures = bundleLecture.objects.filter(bundle=bundle_returned)
		bundle_returned.bundle_lecture_documents = bundleLectureDocument.objects.filter(bundle=bundle_returned)
		bundle_returned.bundle_lecture_slides = bundleLectureSlides.objects.filter(bundle=bundle_returned)

		# for each module in a bundle, attach the module documents and the lectures
		for bundle in bundle_returned.bundle_modules:
			#look up module docs
			moduleDocs = moduleDocuments.objects.filter(module=bundle.module)
			# just get the file name
			for doc in moduleDocs:
				document = str(doc.document)
				document = document.split('/')
				doc.documentName = document[2]

			#attach to the bundle
			bundle.module.moduleDocs = moduleDocs

			#look up the lectures
			moduleLecs = lectures.objects.filter(module=bundle.module).exclude(extracted=False).order_by('title')
			# get the file name
			for lec in moduleLecs:

				# order by lecture number
				first_word = lec.title.strip().lower().split(' ', 1)[0]
				if first_word == 'lecture':
					first_number = lec.title.strip().lower().split(' ')[1].replace('.','').replace(':','')
					try:
						lec.numeric_order = int(first_number)
					except ValueError:
						lec.numeric_order = 0
				else:
					lec.numeric_order = 0

				lecture = str(lec.presentation)
				lecture = lecture.split('/')
				lec.lectureName = lecture[2]
				# get lecture documents
				lectureDocs = lectureDocuments.objects.filter(lecture=lec)
				lec.lectureDocs = lectureDocs
				for lecDoc in lec.lectureDocs:
					document = str(lecDoc.document)
					document = document.split('/')
					lecDoc.documentName = document[2]

			# order titles minus articles
			moduleLecs_ordered = sorted(moduleLecs, key=operator.attrgetter('numeric_order'))

			bundle.module.lectures = moduleLecs_ordered

		#for each lecture, attach lecture docs
		for bundle in bundle_returned.bundle_lectures:
			# get lecture documents
			lectureDocs = lectureDocuments.objects.filter(lecture=bundle.lecture)
			bundle.lecture.lectureDocs = lectureDocs
			for lecDoc in bundle.lecture.lectureDocs:
				document = str(lecDoc.document)
				document = document.split('/')
				lecDoc.documentName = document[2]

		#for each lecture document strip out name of file
		for bundle in bundle_returned.bundle_lecture_documents:
			document = str(bundle.lectureDocument.document)
			document = document.split('/')
			bundle.lectureDocument.documentName = document[2]

		#for each lecture slide strip out name of file and slide notes file
		for bundle in bundle_returned.bundle_lecture_slides:
			slide = str(bundle.lectureSlide.presentation)
			slide = slide.split('/')
			bundle.lectureSlide.slideName = slide[2]


	context_dict = {'bundles_returned':bundles_returned}
	return render(request, 'website/bundle_list.html', context_dict)



@login_required
def updateProfile(request):


	"""
	  Loads the user profile page for editing
	"""


	#get user profile data and pass to view
	try:
		user_profile = profile.objects.get(user=request.user)
	except profile.DoesNotExist:
		user_profile = None

	if request.method == 'POST':
		user_form = UserInfoForm(data=request.POST, instance=request.user)
		profile_form = UserProfileForm(data=request.POST, instance=user_profile)
		if user_form.is_valid() and profile_form.is_valid():
			user_form.save()
			pf = profile_form.save(commit=False)
			try:
				pf.avatar = request.FILES['avatar']
			except KeyError:
				nothing = {}
			try:
				pf.instutution_document = request.FILES['instutution_document']
			except KeyError:
				nothing = {}
			pf.user = request.user
			pf.save()

			return HttpResponseRedirect("/profile/")

		else:
			print user_form.errors, profile_form.errors

	else:
		user_form = UserInfoForm(instance=request.user)
		profile_form = UserProfileForm(instance=user_profile)


	context_dict = {'user_form': user_form, 'profile_form': profile_form}

	return render(request, 'website/update_profile.html', context_dict)


def modulesView(request):

	"""
	  Loads all modules
	"""

	modules_returned = modules.objects.all().order_by('title')

	#remove articles from title for sorting
	for module_returned in modules_returned:
		first_word = module_returned.title.strip().lower().split(' ', 1)[0]
		if first_word == 'a' or first_word == 'the' or first_word == 'and':
			module_returned.no_article_title = module_returned.title.strip().lower().replace(first_word,"",1).strip()
		else:
			module_returned.no_article_title = module_returned.title.strip().lower()

	# order titles minus articles
	modules_returned_ordered = sorted(modules_returned, key=operator.attrgetter('no_article_title'))


	for module_returned in modules_returned_ordered:
		#look up module docs
		moduleDocs = moduleDocuments.objects.filter(module=module_returned).order_by('title')
		# just get the file name
		contents = []
		for doc in moduleDocs:
			document = str(doc.document)
			document = document.split('/')
			doc.documentName = document[2]
			contents.append(doc.document_contents)

		# join document contents together
		module_returned.document_contents = '\n'.join(contents)

		#attach the module docs to the module returned
		module_returned.moduleDocs = moduleDocs

		#look up the lectures
		moduleLecs = lectures.objects.filter(module=module_returned).exclude(extracted=False).order_by('module__title','title')
		# get the file name
		for lec in moduleLecs:
			# order by lecture number
			first_word = lec.title.strip().lower().split(' ', 1)[0]
			if first_word == 'lecture':
				first_number = lec.title.strip().lower().split(' ')[1].replace('.','').replace(':','')
				try:
					lec.numeric_order = int(first_number)
				except ValueError:
					lec.numeric_order = 0
			else:
				lec.numeric_order = 0

			lecture = str(lec.presentation)
			lecture = lecture.split('/')
			lec.lectureName = lecture[2]
			# get lecture documents
			lectureDocs = lectureDocuments.objects.filter(lecture=lec)
			lec.lectureDocs = lectureDocs
			for lecDoc in lec.lectureDocs:
				document = str(lecDoc.document)
				document = document.split('/')
				lecDoc.documentName = document[2]

		# order titles minus articles
		moduleLecs_ordered = sorted(moduleLecs, key=operator.attrgetter('numeric_order'))

		#attach the module docs to the module returned
		module_returned.moduleLecs = moduleLecs_ordered


		# #look up related modules based on keywords in module
		# r_modules_returned = []
		# r_module_documents_returned = []
		# r_lectures_returned = []
		# r_lecture_segments_returned = []
		# r_lecture_documents_returned = []
		# r_lecture_slides_returned = []
		# r_coming_soon_modules_returned = []
		
		# all_results = SearchQuerySet().auto_query(module_returned.description)

		# # sort search query to bins
		# for r in all_results:
		# 	if r.model_name == "modules":
		# 		r_modules_returned.append(r)
		# 	elif r.model_name == "moduledocuments":
		# 		r_module_documents_returned.append(r)
		# 	elif r.model_name == "lectures":
		# 		r_lectures_returned.append(r)

		# # concatonate module querysets
		# # first create list of modules
		# module_modules = []
		# module_documents_modules = []
		# lectures_modules = []
		
		# for module in r_modules_returned:
		# 	if module.object is not None and module.object.title != module_returned.title:
		# 		module_modules.append(module.object)

		# for module_document in r_module_documents_returned:
		# 	if module_document.object is not None and module_document.object.module.title != module_returned.title:
		# 		module_documents_modules.append(module_document.object.module)

		# for lecture in r_lectures_returned:
		# 	if lecture.object is not None and lecture.object.module.title != module_returned.title:
		# 		lectures_modules.append(lecture.object.module)

		# module_list = list(chain(module_modules, module_documents_modules, lectures_modules))

		# # make unique
		# module_set = set(module_list)
		# module_returned.related_modules = list(module_set)

		# pull comments for this module
		module_returned.comments = modulesComments.objects.filter(module=module_returned).order_by('-created')

		for comment in module_returned.comments:
			commenter_profile = profile.objects.get(user=comment.user)
			comment.commenter_profile = commenter_profile

	coming_soon_modules_returned = comingSoonModules.objects.all().order_by('title')

	if request.user.is_authenticated():
		# pull bundles
		bundles_returned = bundles.objects.filter(user=request.user)

	else:
		bundles_returned = bundles.objects.none()

	# comment form
	comment_form = modulesCommentsForm()

	context_dict = {'modules_returned':modules_returned_ordered, 'profile':profile, 'coming_soon_modules_returned':coming_soon_modules_returned, 'bundles_returned':bundles_returned, 'comment_form':comment_form, }
	return render(request, 'website/modules.html', context_dict)

def moduleComments(request, id=None):
	#get module
	module_returned = modules.objects.get(pk=id)

	# pull comments for this module
	module_returned.comments = modulesComments.objects.filter(module=module_returned).order_by('-created')

	for comment in module_returned.comments:
		commenter_profile = profile.objects.get(user=comment.user)
		comment.commenter_profile = commenter_profile

	if request.user.is_authenticated():
		# pull bundles
		bundles_returned = bundles.objects.filter(user=request.user)
	else:
		bundles_returned = bundles.objects.none()

	# comment form
	comment_form = modulesCommentsForm()

	context_dict = {'module_returned':module_returned, 'bundles_returned':bundles_returned, 'comment_form':comment_form}
	return render(request, 'website/module_comments.html', context_dict)


def indexView(request):

	"""
	  Loads all modules
	"""

	modules_returned = modules.objects.all().order_by('title')

	#remove articles from title for sorting
	for module_returned in modules_returned:
		first_word = module_returned.title.strip().lower().split(' ', 1)[0]
		if first_word == 'a' or first_word == 'the' or first_word == 'and':
			module_returned.no_article_title = module_returned.title.strip().lower().replace(first_word,"",1).strip()
		else:
			module_returned.no_article_title = module_returned.title.strip().lower()

	# order titles minus articles
	modules_returned_ordered = sorted(modules_returned, key=operator.attrgetter('no_article_title'))


	for module_returned in modules_returned_ordered:
		#look up module docs
		moduleDocs = moduleDocuments.objects.filter(module=module_returned).order_by('title')
		# just get the file name
		contents = []
		for doc in moduleDocs:
			document = str(doc.document)
			document = document.split('/')
			doc.documentName = document[2]
			contents.append(doc.document_contents)

		# join document contents together
		module_returned.document_contents = '\n'.join(contents)

		#attach the module docs to the module returned
		module_returned.moduleDocs = moduleDocs

		#look up the lectures
		moduleLecs = lectures.objects.filter(module=module_returned).exclude(extracted=False).order_by('module__title','title')
		# get the file name
		for lec in moduleLecs:
			# order by lecture number
			first_word = lec.title.strip().lower().split(' ', 1)[0]
			if first_word == 'lecture':
				first_number = lec.title.strip().lower().split(' ')[1].replace('.','').replace(':','')
				try:
					lec.numeric_order = int(first_number)
				except ValueError:
					lec.numeric_order = 0
			else:
				lec.numeric_order = 0

			lecture = str(lec.presentation)
			lecture = lecture.split('/')
			lec.lectureName = lecture[2]
			# get lecture documents
			lectureDocs = lectureDocuments.objects.filter(lecture=lec)
			lec.lectureDocs = lectureDocs
			for lecDoc in lec.lectureDocs:
				document = str(lecDoc.document)
				document = document.split('/')
				lecDoc.documentName = document[2]

		# order titles minus articles
		moduleLecs_ordered = sorted(moduleLecs, key=operator.attrgetter('numeric_order'))

		#attach the module docs to the module returned
		module_returned.moduleLecs = moduleLecs_ordered


		# pull comments for this module
		module_returned.comments = modulesComments.objects.filter(module=module_returned).order_by('-created')

		for comment in module_returned.comments:
			commenter_profile = profile.objects.get(user=comment.user)
			comment.commenter_profile = commenter_profile

	coming_soon_modules_returned = comingSoonModules.objects.all().order_by('title')

	if request.user.is_authenticated():
		# pull bundles
		bundles_returned = bundles.objects.filter(user=request.user)

	else:
		bundles_returned = bundles.objects.none()

	context_dict = {'modules_returned':modules_returned_ordered, 'coming_soon_modules_returned':coming_soon_modules_returned, 'bundles_returned':bundles_returned }
	return render(request, 'website/index_list.html', context_dict)



def lectureView(request, id=None):

	#get lecture
	lecture_returned = lectures.objects.get(pk=id)

	# get lecture documents
	lectureDocs = lectureDocuments.objects.filter(lecture=lecture_returned)
	lecture_returned.lectureDocs = lectureDocs
	for lecDoc in lecture_returned.lectureDocs:
		document = str(lecDoc.document)
		document = document.split('/')
		lecDoc.documentName = document[2]

	# pull comments for this module
	lecture_returned.comments = lecturesComments.objects.filter(lecture=lecture_returned).order_by('-created')

	for comment in lecture_returned.comments:
		commenter_profile = profile.objects.get(user=comment.user)
		comment.commenter_profile = commenter_profile

	if request.user.is_authenticated():
		# pull bundles
		bundles_returned = bundles.objects.filter(user=request.user)
	else:
		bundles_returned = bundles.objects.none()

	# comment form
	comment_form = lecturesCommentsForm()

	context_dict = {'lecture_returned':lecture_returned, 'bundles_returned':bundles_returned, 'comment_form':comment_form}
	return render(request, 'website/lecture.html', context_dict)

def lectureComments(request, id=None):
	#get lecture
	lecture_returned = lectures.objects.get(pk=id)

	# pull comments for this module
	lecture_returned.comments = lecturesComments.objects.filter(lecture=lecture_returned).order_by('-created')

	for comment in lecture_returned.comments:
		commenter_profile = profile.objects.get(user=comment.user)
		comment.commenter_profile = commenter_profile

	if request.user.is_authenticated():
		# pull bundles
		bundles_returned = bundles.objects.filter(user=request.user)
	else:
		bundles_returned = bundles.objects.none()

	# comment form
	comment_form = lecturesCommentsForm()

	context_dict = {'lecture_returned':lecture_returned, 'bundles_returned':bundles_returned, 'comment_form':comment_form}
	return render(request, 'website/lecture_comments.html', context_dict)



def lecturesView(request):


	"""
	  Loads all lectures grouped by module
	"""

	modules_returned = modules.objects.all().order_by('title')

	#remove articles from title for sorting
	for module_returned in modules_returned:
		first_word = module_returned.title.strip().lower().split(' ', 1)[0]
		if first_word == 'a' or first_word == 'the' or first_word == 'and':
			module_returned.no_article_title = module_returned.title.strip().lower().replace(first_word,"",1).strip()
		else:
			module_returned.no_article_title = module_returned.title.strip().lower()

	# order titles minus articles
	modules_returned_ordered = sorted(modules_returned, key=operator.attrgetter('no_article_title'))


	for module_returned in modules_returned_ordered:
		#look up module docs
		moduleDocs = moduleDocuments.objects.filter(module=module_returned).order_by('title')
		# just get the file name
		contents = []
		for doc in moduleDocs:
			document = str(doc.document)
			document = document.split('/')
			doc.documentName = document[2]
			contents.append(doc.document_contents)

		# join document contents together
		module_returned.document_contents = '\n'.join(contents)

		#attach the module docs to the module returned
		module_returned.moduleDocs = moduleDocs

		#look up the lectures
		moduleLecs = lectures.objects.filter(module=module_returned).exclude(extracted=False).order_by('module__title','title')
		# get the file name
		for lec in moduleLecs:
			# order by lecture number
			first_word = lec.title.strip().lower().split(' ', 1)[0]
			if first_word == 'lecture':
				first_number = lec.title.strip().lower().split(' ')[1].replace('.','').replace(':','')
				try:
					lec.numeric_order = int(first_number)
				except ValueError:
					lec.numeric_order = 0
			else:
				lec.numeric_order = 0

			lecture = str(lec.presentation)
			lecture = lecture.split('/')
			lec.lectureName = lecture[2]
			# get lecture documents
			lectureDocs = lectureDocuments.objects.filter(lecture=lec)
			lec.lectureDocs = lectureDocs
			for lecDoc in lec.lectureDocs:
				document = str(lecDoc.document)
				document = document.split('/')
				lecDoc.documentName = document[2]

			# pull comments for this module
			lec.comments = lecturesComments.objects.filter(lecture=lec).order_by('-created')

			for comment in lec.comments:
				commenter_profile = profile.objects.get(user=comment.user)
				comment.commenter_profile = commenter_profile

		# order titles minus articles
		moduleLecs_ordered = sorted(moduleLecs, key=operator.attrgetter('numeric_order'))

		#attach the module docs to the module returned
		module_returned.moduleLecs = moduleLecs_ordered



	if request.user.is_authenticated():
		# pull bundles
		bundles_returned = bundles.objects.filter(user=request.user)

	else:
		bundles_returned = bundles.objects.none()

	context_dict = {'modules_returned':modules_returned_ordered, 'bundles_returned': bundles_returned }
	return render(request, 'website/lectures.html', context_dict)




@login_required
def membersView(request):
	"""
	  Loads all user profiles
	"""

	profiles_returned = profile.objects.filter(verified=True, public=True).exclude(last_name='', first_name='').order_by(Lower('last_name'), Lower('first_name')).distinct()

	#attach modules and lectures to profiles
	for cp in profiles_returned:
		cp_modules = modules.objects.filter(authors_m2m=cp).order_by('title')
		cp.modules = cp_modules
		cp_csmodules = comingSoonModules.objects.filter(authors_m2m=cp).order_by('title')
		cp.csmodules = cp_csmodules
		#first letter of last name

		cp.first_letter_last_name = cp.last_name[:1]

	context_dict = {'profiles_returned':profiles_returned}
	return render(request, 'website/profiles.html', context_dict)

def searchMembers(request):
	"""
	  AJAX request to search and retrun member profiles
	"""

	if request.method == 'GET':
		#gather variables from get request
		keyword = request.GET.get("keyword","")
		if keyword:
			split_keyword = keyword.split()
			query_str = Q(last_name__icontains=split_keyword[0]) | Q(first_name__icontains=split_keyword[0])
			for kw in split_keyword[1:]:
				query_str.add((Q(last_name__icontains=kw) | Q(first_name__icontains=kw)), query_str.connector)

			profiles_returned = profile.objects.filter(query_str, verified=True, public=True).exclude(last_name='', first_name='').order_by(Lower('last_name'), Lower('first_name')).distinct()
		else:
			profiles_returned = profile.objects.filter(verified=True, public=True).exclude(last_name='', first_name='').order_by(Lower('last_name'), Lower('first_name')).distinct()

		#attach modules and lectures to profiles
		for cp in profiles_returned:
			cp_modules = modules.objects.filter(authors_m2m=cp).order_by('title')
			cp.modules = cp_modules
			cp_csmodules = comingSoonModules.objects.filter(authors_m2m=cp).order_by('title')
			cp.csmodules = cp_csmodules	

			cp.first_letter_last_name = cp.last_name[:1]

	context_dict = {'profiles_returned':profiles_returned}
	return render(request, 'website/profiles_returned.html', context_dict)


def saveSearchString(request):
	"""
	  AJAX request to save search string
	"""

	if request.method == 'GET':
		#gather variables from get request
		searchString = request.GET.get("searchString","")

		# create saved search unless one already exists
		exists = savedSearches.objects.filter(user=request.user, searchString=searchString).count()
		if exists == 0:
			s = savedSearches(user=request.user, searchString=searchString)
			s.save()

	saved_searches = savedSearches.objects.filter(user=request.user)

	context_dict = {'saved_searches':saved_searches}
	return render(request, 'website/saved_searches.html', context_dict)


def saveComment(request):
	"""
	  AJAX request to save comments
	"""

	if request.method == 'GET':
		#gather variables from get request
		itemid = request.GET.get("itemid","")
		comment = request.GET.get("comment","")

		# add item to appropriate bundle
		itemid = itemid.split("_")
		if itemid[0] == 'module':
			#look up module
			module = modules.objects.get(pk=itemid[1])
			#does bundle/module exist already?
			c = modulesComments(user=request.user, module=module, comment=comment)
			c.save()
			return HttpResponseRedirect(reverse('moduleComments', args=(module.id,)))
		elif itemid[0] == 'lecture':
			#look up lecture
			lecture = lectures.objects.get(pk=itemid[1])
			c = lecturesComments(user=request.user, lecture=lecture, comment=comment)
			c.save()
			return HttpResponseRedirect(reverse('lectureComments', args=(lecture.id,)))
		elif itemid[0] == 'lecturesegment':
			#look up lectureSegment
			lectureSegment = lectureSegments.objects.get(pk=itemid[1])
			c = lectureSegmentsComments(user=request.user, lectureSegment=lectureSegment, comment=comment)
			c.save()
			return HttpResponseRedirect(reverse('showLectureSegment', args=(lectureSegment.id,)))
		elif itemid[0] == 'lecturedocument':
			#look up lectureDocument
			lectureDocument = lectureDocuments.objects.get(pk=itemid[1])
			c = lectureDocumentsComments(user=request.user, lectureDocument=lectureDocument, comment=comment)
			c.save()
			return HttpResponseRedirect(reverse('showLectureDocument', args=(lectureDocument.id,)))
		else:
			#look up lectureSlide
			lectureSlide = lectureSlides.objects.get(pk=itemid[1])
			c = lectureSlidesComments(user=request.user, lectureSlide=lectureSlide, comment=comment)
			c.save()
			return HttpResponseRedirect(reverse('showLectureSlide', args=(lectureSlide.id,)))

	return True



@login_required
def dashboard(request):
	"""
	  Check if superuser
	"""
	if request.user.groups.filter(name="superusers").exists():
		empty = {}
	else:
		return HttpResponseRedirect('/')

	"""
	  Admin dashboard -- get modules with their content
	"""

	modulesObjects = modules.objects.all().order_by('title')

	#remove articles from title for sorting
	for module_returned in modulesObjects:
		first_word = module_returned.title.strip().lower().split(' ', 1)[0]
		if first_word == 'a' or first_word == 'the' or first_word == 'and':
			module_returned.no_article_title = module_returned.title.strip().lower().replace(first_word,"",1).strip()
		else:
			module_returned.no_article_title = module_returned.title.strip().lower()

	# order titles minus articles
	modules_returned_ordered = sorted(modulesObjects, key=operator.attrgetter('no_article_title'))


	for module_returned in modules_returned_ordered:
		#look up module docs
		moduleDocs = moduleDocuments.objects.filter(module=module_returned)

		#look up the document name split
		for doc in moduleDocs:
			document = str(doc.document)
			document = document.split('/')
			doc.documentName = document[2]

		#attach the module docs to the module returned
		module_returned.moduleDocumentsObjects = moduleDocs

		#look up the lectures
		moduleLecs = lectures.objects.filter(module=module_returned)

		# now get the lecture documents and get the lecture segments
		for lec in moduleLecs:
			# order by lecture number
			first_word = lec.title.strip().lower().split(' ', 1)[0]
			if first_word == 'lecture':
				first_number = lec.title.strip().lower().split(' ')[1].replace('.','').replace(':','')
				try:
					lec.numeric_order = int(first_number)
				except ValueError:
					lec.numeric_order = 0
				
			else:
				lec.numeric_order = 0

			#look up the document name split
			lecture = str(lec.presentation)
			lecture = lecture.split('/')
			lec.lectureName = lecture[2]

			# get lecture documents
			lectureDocs = lectureDocuments.objects.filter(lecture=lec)
			#look up the document name split
			for doc in lectureDocs:
				document = str(doc.document)
				document = document.split('/')
				doc.documentName = document[2]
			lec.lectureDocumentsObjects = lectureDocs
			lectureSegs = lectureSegments.objects.filter(lecture=lec)
			#look up the document name split
			for lecseg in lectureSegs:
				#look up the document name split
				lectureseg = str(lecseg.presentation)
				lectureseg = lectureseg.split('/')
				lecseg.lectureName = lectureseg[2]
			lec.lectureSegmentsObjects = lectureSegs

		# order titles minus articles
		moduleLecs_ordered = sorted(moduleLecs, key=operator.attrgetter('numeric_order'))

		#attach the module docs to the module returned
		module_returned.lecturesObjects = moduleLecs_ordered


	context_dict = {'modulesObjects': modules_returned_ordered}
	return render(request, 'website/dashboard.html', context_dict)


# admin form views
@login_required
def admin_module(request, id=None):
	"""
	  Check if superuser
	"""
	if request.user.groups.filter(name="superusers").exists():
		empty = {}
	else:
		return HttpResponseRedirect('/')

	if id:
		modulesObject = modules.objects.get(id=id)
	else:
		modulesObject = modules()

	# A HTTP POST?
	if request.method == 'POST':
		form = modulesForm(request.POST, request.FILES, instance=modulesObject)

		# Have we been provided with a valid form?
		if form.is_valid():
			# Save the new data to the database.
			f = form.save(commit=False)
			f.save()
			# Without this next line the tags won't be saved.
			form.save_m2m()

			# ensure each author is tagged as contributing
			profile_objects = profile.objects.filter(pk=modulesObject.authors_m2m.all())
			for profile_object in profile_objects:
				profile_object.contributing = True
				profile_object.save()


			# route user depending on what button they clicked
			if 'save' in request.POST:
				# send back to dashboard
				return HttpResponseRedirect('/dashboard/')
			elif 'new_module_document' in request.POST:
				# send to new module doc form
				return HttpResponseRedirect(reverse('admin_moduledoc', args=(0,f.id)))
			else:
				# send to new lecture form
				return HttpResponseRedirect(reverse('admin_lecture', args=(0,f.id)))

		else:
			# The supplied form contained errors - just print them to the terminal.
			print form.errors
	else:
		# If the request was not a POST, display the form to enter details.
		form = modulesForm(instance=modulesObject)
		form.fields["authors_m2m"].queryset = profile.objects.filter(verified=True).exclude(last_name='', first_name='').order_by('last_name', 'first_name')

	# Bad form (or form details), no form supplied...
	# Render the form with error messages (if any).
	return render(request, 'website/admin_modules.html', {'form': form, 'media': form.media})


@login_required
def admin_removemodule(request, id=None):
	"""
	  Check if superuser
	"""
	if request.user.groups.filter(name="superusers").exists():
		empty = {}
	else:
		return HttpResponseRedirect('/')

	if id:
		modulesObject = modules.objects.get(id=id)
			#look up module docs
		moduleDocs = moduleDocuments.objects.filter(module=modulesObject)
		# just get the file name
		for doc in moduleDocs:
			document = str(doc.document)
			document = document.split('/')
			doc.documentName = document[2]

		#look up the lectures
		moduleLecs = lectures.objects.filter(module=modulesObject).order_by('title')
		# get the file name
		for lec in moduleLecs:
			lecture = str(lec.presentation)
			lecture = lecture.split('/')
			lec.lectureName = lecture[2]
			# get lecture documents
			lectureDocs = lectureDocuments.objects.filter(lecture=lec)
			lec.lectureDocs = lectureDocs
			for lecDoc in lec.lectureDocs:
				document = str(lecDoc.document)
				document = document.split('/')
				lecDoc.documentName = document[2]

	else:
		return HttpResponseRedirect('/dashboard/')

	# A HTTP POST?
	if request.method == 'POST':
		form = modulesRemoveForm(request.POST, instance=modulesObject)

		# Have we been provided with a valid form?
		if form.is_valid():
			# if form submitted, delete module
			if 'delete' in request.POST:
				modulesObject.delete()
				return HttpResponseRedirect(request.POST['referer'])

		else:
			# The supplied form contained errors - just print them to the terminal.
			print form.errors
	else:
		# If the request was not a POST, display the form to enter details.
		form = modulesRemoveForm(instance=modulesObject)

	# Bad form (or form details), no form supplied...
	# Render the form with error messages (if any).
	return render(request, 'website/admin_removemodules.html', {'form': form, 'modulesObject': modulesObject, 'moduleDocs': moduleDocs, 'moduleLecs': moduleLecs})



@login_required
def admin_moduledoc(request, id=None, moduleid=None):
	"""
	  Check if superuser
	"""
	if request.user.groups.filter(name="superusers").exists():
		empty = {}
	else:
		return HttpResponseRedirect('/')

	if id and int(id) > 0:
		moduleDocumentsObject = moduleDocuments.objects.get(id=id)
	else:
		moduleDocumentsObject = moduleDocuments()

	# A HTTP POST?
	if request.method == 'POST':
		form = moduleDocumentsForm(request.POST, request.FILES, instance=moduleDocumentsObject)

		# Have we been provided with a valid form?
		if form.is_valid():
			# Save the new data to the database.
			f = form.save(commit=False)
			f.save()
			# Without this next line the tags won't be saved.
			form.save_m2m()

			# route user depending on what button they clicked
			if 'save' in request.POST:
				# send back to dashboard
				return HttpResponseRedirect('/dashboard/')
			else:
				# send to new module doc form
				return HttpResponseRedirect(reverse('admin_moduledoc', args=(0,f.module.id,)))

		else:
			# The supplied form contained errors - just print them to the terminal.
			print form.errors
	else:
		# If the request was not a POST, display the form to enter details.
		form = moduleDocumentsForm(instance=moduleDocumentsObject)

	# Bad form (or form details), no form supplied...
	# Render the form with error messages (if any).
	return render(request, 'website/admin_moduledoc.html', {'form': form, 'media': form.media, 'moduleid': moduleid})


@login_required
def admin_removemoduledoc(request, id=None):
	"""
	  Check if superuser
	"""
	if request.user.groups.filter(name="superusers").exists():
		empty = {}
	else:
		return HttpResponseRedirect('/')

	if id:
		moduleDocumentsObject = moduleDocuments.objects.get(id=id)
		document = str(moduleDocumentsObject.document)
		document = document.split('/')
		moduleDocumentsObject.document.documentName = document[2]

	# A HTTP POST?
	if request.method == 'POST':
		form = moduleDocumentsRemoveForm(request.POST, instance=moduleDocumentsObject)

		# Have we been provided with a valid form?
		if form.is_valid():
			# if form submitted, delete module
			if 'delete' in request.POST:
				moduleDocumentsObject.delete()
				return HttpResponseRedirect(request.POST['referer'])

		else:
			# The supplied form contained errors - just print them to the terminal.
			print form.errors
	else:
		# If the request was not a POST, display the form to enter details.
		form = moduleDocumentsRemoveForm(instance=moduleDocumentsObject)

	# Bad form (or form details), no form supplied...
	# Render the form with error messages (if any).
	return render(request, 'website/admin_removemoduledoc.html', {'form': form, 'moduleDocumentsObject': moduleDocumentsObject})



@login_required
def admin_lecture(request, id=None, moduleid=None):
	"""
	  Check if superuser
	"""
	if request.user.groups.filter(name="superusers").exists():
		empty = {}
	else:
		return HttpResponseRedirect('/')

	if id and int(id) > 0:
		lectureObject = lectures.objects.get(id=id)
	else:
		lectureObject = lectures()

	# A HTTP POST?
	if request.method == 'POST':
		form = lectureForm(request.POST, request.FILES, instance=lectureObject)

		# Have we been provided with a valid form?
		if form.is_valid():
			# Save the new data to the database.
			f = form.save(commit=False)

			# check to see if the field "presentation" has not changed and if so set extracted to True, otherwise set to false
			if 'presentation' in form.changed_data:
				f.extracted = False
				# remove any associated lecture slides
				slides = lectureSlides.objects.filter(lecture=lectureObject)
				for slide in slides:
					slide.delete()
			else:
				f.extracted = True

			f.save()

			# Without this next line the tags won't be saved.
			form.save_m2m()

			# flag any associated lecture segments for review
			lectureSegs = lectureSegments.objects.filter(lecture=lectureObject)
			for seg in lectureSegs:
				seg.updated_lecture_review = True
				seg.save()

			# route user depending on what button they clicked
			if 'save' in request.POST:
				# send back to dashboard
				return HttpResponseRedirect('/dashboard/')
			elif 'new_lecure_document' in request.POST:
				# send to new lecture doc form
				return HttpResponseRedirect(reverse('admin_lecturedoc', args=(0,f.id)))
			elif 'new_lecure_segment' in request.POST:
				# send to new lecture doc form
				return HttpResponseRedirect(reverse('admin_lecturesegment', args=(0,f.id)))
			else:
				# send to new module doc form
				return HttpResponseRedirect(reverse('admin_lecture', args=(0,f.module.id,)))

		else:
			# The supplied form contained errors - just print them to the terminal.
			print form.errors
	else:
		# If the request was not a POST, display the form to enter details.
		form = lectureForm(instance=lectureObject)
		form.fields["authors_m2m"].queryset = profile.objects.filter(verified=True).exclude(last_name='', first_name='').order_by('last_name', 'first_name')

	# Bad form (or form details), no form supplied...
	# Render the form with error messages (if any).
	return render(request, 'website/admin_lecture.html', {'form': form, 'media': form.media, 'moduleid': moduleid})


@login_required
def admin_removelecture(request, id=None):
	"""
	  Check if superuser
	"""
	if request.user.groups.filter(name="superusers").exists():
		empty = {}
	else:
		return HttpResponseRedirect('/')

	if id:
		lectureObject = lectures.objects.get(id=id)
		document = str(lectureObject.presentation)
		document = document.split('/')
		lectureObject.presentation.documentName = document[2]

		# get lecture documents
		lectureDocs = lectureDocuments.objects.filter(lecture=lectureObject)
		for lecDoc in lectureDocs:
			document = str(lecDoc.document)
			document = document.split('/')
			lecDoc.documentName = document[2]

		# get lecture segments
		lectureSegs = lectureSegments.objects.filter(lecture=lectureObject)
		for seg in lectureSegs:
			document = str(seg.presentation)
			document = document.split('/')
			seg.documentName = document[2]


	# A HTTP POST?
	if request.method == 'POST':
		form = lectureRemoveForm(request.POST, instance=lectureObject)

		# Have we been provided with a valid form?
		if form.is_valid():
			# if form submitted, delete module
			if 'delete' in request.POST:
				lectureObject.delete()
				return HttpResponseRedirect(request.POST['referer'])

		else:
			# The supplied form contained errors - just print them to the terminal.
			print form.errors
	else:
		# If the request was not a POST, display the form to enter details.
		form = lectureRemoveForm(instance=lectureObject)

	# Bad form (or form details), no form supplied...
	# Render the form with error messages (if any).
	return render(request, 'website/admin_removelecture.html', {'form': form, 'lectureObject': lectureObject, 'lectureDocs': lectureDocs, 'lectureSegs': lectureSegs})



@login_required
def admin_lecturesegment(request, id=None, lectureid=None):
	"""
	  Check if superuser
	"""
	if request.user.groups.filter(name="superusers").exists():
		empty = {}
	else:
		return HttpResponseRedirect('/')

	if id and int(id) > 0:
		lecturesegmentObject = lectureSegments.objects.get(id=id)
	else:
		lecturesegmentObject = lectureSegments()

	# A HTTP POST?
	if request.method == 'POST':
		form = lecturesegmentForm(request.POST, instance=lecturesegmentObject)

		# Have we been provided with a valid form?
		if form.is_valid():
			# Save the new data to the database.
			f = form.save(commit=False)
			#ensure extracted is False and lecture review is False -- important after lecture has been updated
			f.extracted = False
			f.updated_lecture_review = False

			# create a ppt file using the slide numbers
			# slides are 0-indexed so the number the admin enters are one larger than the slide numbers used to manipulate the slides
			prs = Presentation(f.lecture.presentation)

			# calculate the last slide number
			largestslidenumber = len(prs.slides)

			# loop through a range from the maximum slide number plus 1 (which equals f.maxslidenumber) to the largest slide number
			for index in xrange(f.maxslidenumber, largestslidenumber):
				prs.slides.delete_slide(prs, f.maxslidenumber)

			minslide = f.minslidenumber - 1
			for index in xrange(0, minslide):
				prs.slides.delete_slide(prs, 0)


			path_to_file = MEDIA_ROOT + '/' + str(f.lecture.presentation) + "_" + str(f.minslidenumber) + "-" + str(f.maxslidenumber) + ".pptx"
			prs.save(path_to_file)

			seg = open(path_to_file)
			seg_file = File(seg)
			head, tail = os.path.split(path_to_file)
			f.presentation.save(tail, seg_file)
			os.remove(path_to_file)

			f.save()


			# route user depending on what button they clicked
			if 'save' in request.POST:
				# send back to dashboard
				return HttpResponseRedirect('/dashboard/')
			else:
				# send to new module doc form
				return HttpResponseRedirect(reverse('admin_lecturesegment', args=(0,f.lecture.id,)))

		else:
			# The supplied form contained errors - just print them to the terminal.
			print form.errors
	else:
		# If the request was not a POST, display the form to enter details.
		form = lecturesegmentForm(instance=lecturesegmentObject)

	# Bad form (or form details), no form supplied...
	# Render the form with error messages (if any).
	return render(request, 'website/admin_lecturesegment.html', {'form': form, 'media': form.media, 'lectureid': lectureid})


@login_required
def admin_removelecturesegment(request, id=None):
	"""
	  Check if superuser
	"""
	if request.user.groups.filter(name="superusers").exists():
		empty = {}
	else:
		return HttpResponseRedirect('/')

	if id:
		lecturesegmentObject = lectureSegments.objects.get(id=id)
		document = str(lecturesegmentObject.presentation)
		document = document.split('/')
		lecturesegmentObject.presentation.documentName = document[2]

	# A HTTP POST?
	if request.method == 'POST':
		form = lecturesegmentRemoveForm(request.POST, instance=lecturesegmentObject)

		# Have we been provided with a valid form?
		if form.is_valid():
			# if form submitted, delete module
			if 'delete' in request.POST:
				lecturesegmentObject.delete()
				return HttpResponseRedirect(request.POST['referer'])

		else:
			# The supplied form contained errors - just print them to the terminal.
			print form.errors
	else:
		# If the request was not a POST, display the form to enter details.
		form = lecturesegmentRemoveForm(instance=lecturesegmentObject)

	# Bad form (or form details), no form supplied...
	# Render the form with error messages (if any).
	return render(request, 'website/admin_removelecturesegment.html', {'form': form, 'lecturesegmentObject': lecturesegmentObject})



@login_required
def admin_lecturedoc(request, id=None, lectureid=None):
	"""
	  Check if superuser
	"""
	if request.user.groups.filter(name="superusers").exists():
		empty = {}
	else:
		return HttpResponseRedirect('/')

	if id and int(id) > 0:
		lecturedocObject = lectureDocuments.objects.get(id=id)
	else:
		lecturedocObject = lectureDocuments()

	# A HTTP POST?
	if request.method == 'POST':
		form = lecturedocForm(request.POST, request.FILES, instance=lecturedocObject)

		# Have we been provided with a valid form?
		if form.is_valid():
			# Save the new data to the database.
			f = form.save(commit=False)
			f.save()
			# Without this next line the tags won't be saved.
			form.save_m2m()

			# route user depending on what button they clicked
			if 'save' in request.POST:
				# send back to dashboard
				return HttpResponseRedirect('/dashboard/')
			else:
				# send to new lecture doc form
				return HttpResponseRedirect(reverse('admin_lecturedoc', args=(0,f.lecture.id,)))

		else:
			# The supplied form contained errors - just print them to the terminal.
			print form.errors
	else:
		# If the request was not a POST, display the form to enter details.
		form = lecturedocForm(instance=lecturedocObject)

	# Bad form (or form details), no form supplied...
	# Render the form with error messages (if any).
	return render(request, 'website/admin_lecturedoc.html', {'form': form, 'media': form.media, 'lectureid': lectureid})


@login_required
def admin_removelecturedoc(request, id=None):
	"""
	  Check if superuser
	"""
	if request.user.groups.filter(name="superusers").exists():
		empty = {}
	else:
		return HttpResponseRedirect('/')

	if id:
		lecturedocObject = lectureDocuments.objects.get(id=id)
		document = str(lecturedocObject.document)
		document = document.split('/')
		lecturedocObject.document.documentName = document[2]

	# A HTTP POST?
	if request.method == 'POST':
		form = lecturedocRemoveForm(request.POST, instance=lecturedocObject)

		# Have we been provided with a valid form?
		if form.is_valid():
			# if form submitted, delete module
			if 'delete' in request.POST:
				lecturedocObject.delete()
				return HttpResponseRedirect(request.POST['referer'])

		else:
			# The supplied form contained errors - just print them to the terminal.
			print form.errors
	else:
		# If the request was not a POST, display the form to enter details.
		form = lecturedocRemoveForm(instance=lecturedocObject)

	# Bad form (or form details), no form supplied...
	# Render the form with error messages (if any).
	return render(request, 'website/admin_removelecturedoc.html', {'form': form, 'lecturedocObject': lecturedocObject})


@login_required
def admin_accounts(request):
	"""
	  Check if superuser
	"""
	if request.user.groups.filter(name="superusers").exists():
		empty = {}
	else:
		return HttpResponseRedirect('/')

	# unverified
	unverified = profile.objects.filter(verified__exact=None).exclude(last_name='', first_name='').order_by('name')

	# accepted
	accepted = profile.objects.filter(verified__exact=True).exclude(last_name='', first_name='').order_by('name')

	# rejected
	rejected = profile.objects.filter(verified__exact=False).exclude(last_name='', first_name='').order_by('name')

	# pull all account holders
	return render(request, 'website/admin_accounts.html', {'unverified': unverified, 'accepted': accepted, 'rejected':rejected})

@login_required
def admin_unverified(request):
	"""
	  Check if superuser
	"""
	if request.user.groups.filter(name="superusers").exists():
		empty = {}
	else:
		return HttpResponseRedirect('/')

	# unverified
	unverified = profile.objects.filter(verified__exact=None).order_by('name')

	# pull all account holders
	return render(request, 'website/admin_unverified.html', {'unverified': unverified,})

@login_required
def admin_update_profile(request, id=None):
	"""
	  Check if superuser
	"""
	if request.user.groups.filter(name="superusers").exists():
		empty = {}
	else:
		return HttpResponseRedirect('/')

	#get user profile data and pass to view
	thisUser = User.objects.get(pk=id)
	try:
		user_profile = profile.objects.get(user=thisUser)
	except profile.DoesNotExist:
		user_profile = None

	if request.method == 'POST':
		user_form = UserInfoForm(data=request.POST, instance=thisUser)
		profile_form = AdminUserProfileForm(data=request.POST, instance=user_profile)
		if user_form.is_valid() and profile_form.is_valid():
			user_form.save()
			pf = profile_form.save(commit=False)
			try:
				pf.avatar = request.FILES['avatar']
			except KeyError:
				nothing = {}
			try:
				pf.instutution_document = request.FILES['instutution_document']
			except KeyError:
				nothing = {}
			pf.save()

			return HttpResponseRedirect(request.POST['referer'])

		else:
			print user_form.errors, profile_form.errors

	else:
		user_form = UserInfoForm(instance=thisUser)
		profile_form = AdminUserProfileForm(instance=user_profile)


	context_dict = {'user_form': user_form, 'profile_form': profile_form}

	return render(request, 'website/admin_update_profile.html', context_dict)

@login_required
def admin_verify_user(request, id=None):
	"""
	  Check if superuser
	"""
	if request.user.groups.filter(name="superusers").exists():
		empty = {}
	else:
		return HttpResponseRedirect('/')

	#get user profile data and pass to view
	thisUser = User.objects.get(pk=id)
	try:
		user_profile = profile.objects.get(user=thisUser)
	except profile.DoesNotExist:
		user_profile = None

	if request.method == 'POST':
		verify_form = AdminVerifyUserForm(data=request.POST, instance=user_profile)
		if verify_form.is_valid():
			vf = verify_form.save()

			#send email if use was verified or rejected
			if vf.verified != None:
				if vf.verified:
					#send email
					subject = "[GAHTC] Welcome to GAHTC!"
					html_message = "<p><h1>Welcome to GAHTC!</h1><p>Thank you for your interest in the Global Architectural History Collaborative, generously supported by the Andrew W. Mellon foundation and in connection with MIT and the Society of Architectural Historians (SAH).</p><p>The ambition of the GAHTC is to address the needs of educators in diverse disciplinary contexts by providing practical lecture materials for teaching global architectural history at the survey level. This effort does not preclude more advanced level education, but the main purpose of the Collaborative is to transform the discipline 'from below'-that is, to help shape the discourse of architectural history by reshaping its teaching at the survey level. GAHTC's current online teaching library of nearly 200 lectures are all produced to strengthen and elevate the global perspective within programs teaching architectural history. The library is continuing to grow with relevant content.</p><p>The GAHTC will dedicate its newest grant of a $1.5 million to promoting the development of survey course material in the history of architecture, thus strengthening its position within humanities teaching, while also sponsoring teacher-to-teacher conversations that support pedagogy with a global perspective.</p><p>To accomplish these goals, the GAHTC has created six new funding opportunities for research and teaching. You can find more information on the various grants <a href='http://gahtc.org/pages/grants/'>here</a>.</p><p>Guidelines for submission as well as prospective grantee resources can be found <a href='http://gahtc.org/index/'>here</a>, under the 'Resources' heading.</p><p>Please direct all questions related to membership and grant applications to our Project Manager, Eliana Hamdi Murchie at emurchie@mit.edu.</p><p>We look forward to your contribution to the GAHTC community!</p></p>"
					message = "<p><h1>Welcome to GAHTC!</h1><p>Thank you for your interest in the Global Architectural History Collaborative, generously supported by the Andrew W. Mellon foundation and in connection with MIT and the Society of Architectural Historians (SAH).</p><p>The ambition of the GAHTC is to address the needs of educators in diverse disciplinary contexts by providing practical lecture materials for teaching global architectural history at the survey level. This effort does not preclude more advanced level education, but the main purpose of the Collaborative is to transform the discipline 'from below'-that is, to help shape the discourse of architectural history by reshaping its teaching at the survey level. GAHTC's current online teaching library of nearly 200 lectures are all produced to strengthen and elevate the global perspective within programs teaching architectural history. The library is continuing to grow with relevant content.</p><p>The GAHTC will dedicate its newest grant of a $1.5 million to promoting the development of survey course material in the history of architecture, thus strengthening its position within humanities teaching, while also sponsoring teacher-to-teacher conversations that support pedagogy with a global perspective.</p><p>To accomplish these goals, the GAHTC has created six new funding opportunities for research and teaching. You can find more information on the various grants <a href='http://gahtc.org/pages/grants/'>here</a>.</p><p>Guidelines for submission as well as prospective grantee resources can be found <a href='http://gahtc.org/index/'>here</a>, under the 'Resources' heading.</p><p>Please direct all questions related to membership and grant applications to our Project Manager, Eliana Hamdi Murchie at emurchie@mit.edu.</p><p>We look forward to your contribution to the GAHTC community!</p></p>"

					send_mail(subject, message, 'gahtcweb@gmail.com', [user_profile.user.email], fail_silently=True, html_message=html_message)
				else:
					subject = "[GAHTC] We're sorry, but your account has been rejected."
					html_message = "Hello "+ user_profile.name +",<br /><br />We're sorry to report, but your account had been rejected by the GAHTC administrator. You may still search the GAHTC site for course materials, but you will not be able to download any of these materials until your account is approved. If you feel that this has been done in error, please contact GAHTC by visiting <a href='http://gahtc.org/contact/'>Contact GAHTC</a>.<br /><br />GAHTC | Global Architectural History Teaching Collaborative"
					message = "Hello "+ user_profile.name +", We're sorry to report, but your account had been rejected by the GAHTC administrator. You may still search the GAHTC site for course materials, but you will not be able to download any of these materials until your account is approved. If you feel that this has been done in error, please contact GAHTC by visiting the Contact GAHTC page ( http://gahtc.org/contact/ ). GAHTC | Global Architectural History Teaching Collaborative"

					send_mail(subject, message, 'gahtcweb@gmail.com', [user_profile.user.email], fail_silently=True, html_message=html_message)


			return HttpResponseRedirect(request.POST['referer'])

		else:
			print verify_form.errors

	else:
		verify_form = AdminVerifyUserForm(instance=user_profile)


	context_dict = {'verify_form': verify_form, 'user_profile': user_profile}

	return render(request, 'website/admin_verify_user.html', context_dict)



@login_required
def admin_download_user_profiles(request):
	"""
	  Check if superuser
	"""
	if request.user.groups.filter(name="superusers").exists():
		empty = {}
	else:
		return HttpResponseRedirect('/')

	#loop thorough users to create all other rows
	profiles = profile.objects.values('user__username', 'name', 'user__email', 'user__last_login', 'user__date_joined', 'institution', 'institution_address', 'institution_city', 'institution_country', 'institution_postal_code', 'teaching', 'introduction', 'avatar', 'title', 'website', 'instutution_document', 'verified')

	return render_to_csv_response(profiles)


@login_required
def admin_downloads(request):
	"""
	  Check if superuser
	"""
	if request.user.groups.filter(name="superusers").exists():
		empty = {}
	else:
		return HttpResponseRedirect('/')

	bundles_downloaded = bundles.objects.filter(downloaded=True).order_by('-created')
	for b in bundles_downloaded:
		# look up the materials in each bundle
		b.bundleModules = bundleModule.objects.filter(bundle__exact=b)
		b.bundleLectures = bundleLecture.objects.filter(bundle__exact=b)
		b.bundleLectureSegments = bundleLectureSegments.objects.filter(bundle__exact=b)
		b.bundleLectureDocuments = bundleLectureDocument.objects.filter(bundle__exact=b)
		b.bundleLectureSlides = bundleLectureSlides.objects.filter(bundle__exact=b)

	modules_downloaded = userModuleDownload.objects.filter(downloaded=True).order_by('-created')
	lectures_downloaded = userLectureDownload.objects.filter(downloaded=True).order_by('-created')

	context_dict = {'bundles_downloaded':bundles_downloaded,'modules_downloaded':modules_downloaded,'lectures_downloaded':lectures_downloaded}

	return render(request, 'website/admin_downloads.html', context_dict)

@login_required
def admin_coming_soon_module(request, id=None):
	"""
	  Check if superuser
	"""
	if request.user.groups.filter(name="superusers").exists():
		empty = {}
	else:
		return HttpResponseRedirect('/')

	if id:
		modulesObject = comingSoonModules.objects.get(id=id)
	else:
		modulesObject = comingSoonModules()

	# A HTTP POST?
	if request.method == 'POST':
		form = CSmodulesForm(request.POST, request.FILES, instance=modulesObject)

		# Have we been provided with a valid form?
		if form.is_valid():
			# Save the new data to the database.
			f = form.save()

			# send back to dashboard
			return HttpResponseRedirect('/dashboard/')
			
		else:
			# The supplied form contained errors - just print them to the terminal.
			print form.errors
	else:
		# If the request was not a POST, display the form to enter details.
		form = CSmodulesForm(instance=modulesObject)
		form.fields["authors_m2m"].queryset = profile.objects.filter(verified=True).exclude(last_name='', first_name='').order_by('last_name', 'first_name')

	return render(request, 'website/admin_coming_soon_module.html', {'form': form, 'media': form.media})


@login_required
def admin_remove_coming_soon_modules(request):
	"""
	  Check if superuser
	"""
	if request.user.groups.filter(name="superusers").exists():
		empty = {}
	else:
		return HttpResponseRedirect('/')

	csModules = comingSoonModules.objects.all()
	return render(request, 'website/admin_remove_coming_soon_modules.html', {'csModules': csModules,})

@login_required
def admin_remove_coming_soon_module(request, id=None):
	"""
	  Check if superuser
	"""
	if request.user.groups.filter(name="superusers").exists():
		empty = {}
	else:
		return HttpResponseRedirect('/')

	if id:
		modulesObject = comingSoonModules.objects.get(id=id)
	else:
		modulesObject = comingSoonModules()

	# A HTTP POST?
	if request.method == 'POST':
		form = CSmodulesRemoveForm(request.POST, instance=modulesObject)

		# Have we been provided with a valid form?
		if form.is_valid():
			# if form submitted, delete module
			if 'delete' in request.POST:
				modulesObject.delete()
				return HttpResponseRedirect(request.POST['referer'])

		else:
			# The supplied form contained errors - just print them to the terminal.
			print form.errors
	else:
		# If the request was not a POST, display the form to enter details.
		form = CSmodulesRemoveForm(instance=modulesObject)

	# Bad form (or form details), no form supplied...
	# Render the form with error messages (if any).
	return render(request, 'website/admin_remove_coming_soon_module.html', {'form': form, 'modulesObject': modulesObject,})

@login_required
def admin_managedoctypes(request, id=None):
	"""
	  Check if superuser
	"""
	if request.user.groups.filter(name="superusers").exists():
		empty = {}
	else:
		return HttpResponseRedirect('/')

	#get doc types
	doc_types = docType.objects.all()

	if id:
		doc_type_object = docType.objects.get(id=id)
	else:
		doc_type_object = docType()

	# A HTTP POST?
	if request.method == 'POST':
		form = docTypeAddForm(request.POST, instance=doc_type_object)

		# Have we been provided with a valid form?
		if form.is_valid():
			# Save the new data to the database.
			f = form.save()

			# send back to dashboard
			return HttpResponseRedirect('/dashboard/')

		else:
			# The supplied form contained errors - just print them to the terminal.
			print form.errors
	else:
		# If the request was not a POST, display the form to enter details.
		form = docTypeAddForm(instance=doc_type_object)

	return render(request, 'website/admin_managedoctypes.html', {'doc_types': doc_types, 'form': form})

def admin_removedoctype(request):

	"""
	  AJAX request to save search string
	"""

	if request.method == 'GET':
		#gather variables from get request
		doctypeid = request.GET.get("doctypeid","")

		# look up lecture and module docs that have this doc type and set them to nno doc type so they don't get deleted!
		doc_type_object = docType.objects.get(pk=doctypeid)
		lec_docs = lectureDocuments.objects.filter(doc_type=doc_type_object)

		for lec_doc in lec_docs:
			print lec_doc
			lec_doc.doc_type = None

		# delete document type
		doc_type_object.delete()


	#get doc types
	doc_types = docType.objects.all()


	doc_type_object = docType()

	# If the request was not a POST, display the form to enter details.
	form = docTypeAddForm(instance=doc_type_object)


	context_dict = {'doc_types':doc_types, 'form': form}
	return render(request, 'website/admin_current_doc_types.html', context_dict)


def contactBundle(request):
	"""
	  AJAX request to save if user wants to be contaced about their bundle
	"""

	if request.method == 'GET':
		#gather variables from get request
		bundleid = request.GET.get("bundleid","")

		#look up bundle
		bundle = bundles.objects.get(pk=bundleid)
		bundle.contact = True
		bundle.save()

	return JsonResponse({'foo': 'bar'})

def dontContactBundle(request):
	"""
	  AJAX request to save if user wants to be contaced about their bundle
	"""

	if request.method == 'GET':
		#gather variables from get request
		bundleid = request.GET.get("bundleid","")

		#look up bundle
		bundle = bundles.objects.get(pk=bundleid)
		bundle.contact = False
		bundle.save()

	return JsonResponse({'foo': 'bar'})

def contactModule(request):
	"""
	  AJAX request to save if user wants to be contaced about their module
	"""

	if request.method == 'GET':
		#gather variables from get request
		moduleid = request.GET.get("moduleid","")

		#look up userModuleDownload
		c = userModuleDownload.objects.filter(module__exact=moduleid, user__exact=request.user)
		for o in c:
			o.contact = True
			o.save()

	return JsonResponse({'foo': 'bar'})

def dontContactModule(request):
	"""
	  AJAX request to save if user wants to be contaced about their module
	"""

	if request.method == 'GET':
		#gather variables from get request
		moduleid = request.GET.get("moduleid","")

		#look up userModuleDownload
		c = userModuleDownload.objects.filter(module__exact=moduleid, user__exact=request.user)
		for o in c:
			o.contact = False
			o.save()

	return JsonResponse({'foo': 'bar'})

def contactLecture(request):
	"""
	  AJAX request to save if user wants to be contaced about their lecture
	"""

	if request.method == 'GET':
		#gather variables from get request
		lectureid = request.GET.get("lectureid","")

		#look up userLectureDownload
		c = userLectureDownload.objects.filter(lecture__exact=lectureid, user__exact=request.user)
		for o in c:
			o.contact = True
			o.save()

	return JsonResponse({'foo': 'bar'})

def dontContactLecture(request):
	"""
	  AJAX request to save if user wants to be contaced about their lecture
	"""

	if request.method == 'GET':
		#gather variables from get request
		lectureid = request.GET.get("lectureid","")

		#look up userModuleDownload
		c = userLectureDownload.objects.filter(lecture__exact=lectureid, user__exact=request.user)
		for o in c:
			o.contact = False
			o.save()

	return JsonResponse({'foo': 'bar'})



def inlineEditUpdateTitle(request):
	"""
	  AJAX request to save new title for an item from the admin dashboard
	"""

	if request.method == 'GET':
		#gather variables from get request
		#/inline_edit_update_title/?new_text=" + new_text + "&type=" + type + "&itemid=" + itemid,
		new_text = request.GET.get("new_text","")
		model_type = request.GET.get("type","")
		itemid = request.GET.get("itemid","")

		# check type
		if model_type == "moduledoc":
			obj = moduleDocuments.objects.get(pk=itemid)
		elif model_type == "lecture":
			obj = lectures.objects.get(pk=itemid)
		elif model_type == "lecturesegment":
			obj = lectureSegments.objects.get(pk=itemid)
		elif model_type == "lecturedocument":
			obj = lectureDocuments.objects.get(pk=itemid)

		obj.title = new_text
		obj.save()

	return JsonResponse({'foo': 'bar'})



def inlineEditUpdateDoc(request):
	"""
	  AJAX request to save new title for an item from the admin dashboard
	"""

	if request.method == 'GET':
		#gather variables from get request
		#/inline_edit_update_title/?new_text=" + new_text + "&type=" + type + "&itemid=" + itemid,
		new_text = request.GET.get("new_text","")
		model_type = request.GET.get("type","")
		itemid = request.GET.get("itemid","")

		# check type
		if model_type == "moduledoc":
			obj = moduleDocuments.objects.get(pk=itemid)
			doc = obj.document
		elif model_type == "lecture":
			obj = lectures.objects.get(pk=itemid)
			doc = obj.presentation
		elif model_type == "lecturesegment":
			obj = lectureSegments.objects.get(pk=itemid)
			doc = obj.presentation
		elif model_type == "lecturedocument":
			obj = lectureDocuments.objects.get(pk=itemid)
			doc = obj.document

		# rename document on file server
		head, tail = os.path.split(doc.path)
		new_path = head + '/' + new_text

		#change name
		doc_name = doc.name.split('/')

		#update fileserver name
		os.rename(doc.path, new_path)

		#save to database
		doc.name = doc_name[0] + '/' + doc_name[1] + '/' + new_text
		obj.save()


	return JsonResponse({'foo': 'bar'})
