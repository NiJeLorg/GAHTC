import sys,os
from django.core.management.base import BaseCommand, CommandError
from django.conf import settings
from django.core.files import File
from subprocess import call
from website.models import *
import pptx
from pptx import Presentation
import string
from textblob import TextBlob
import gc

"""
  Opens pptx files and extracts all text
  Requires libreoffice and Imagemagick
"""

from pptx.util import lazyproperty, Pt
from pptx.parts.slide import BaseSlide, Slide, _SlideShapeTree, _SlidePlaceholders
from pptx.shapes.shape import BaseShape
from pptx.opc.constants import RELATIONSHIP_TYPE as RT, CONTENT_TYPE as CT
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER_TYPE
from pptx.oxml.xmlchemy import BaseOxmlElement, OneAndOnlyOne, ZeroOrOne
from pptx.oxml import parse_xml
from pptx.oxml.ns import nsdecls
from pptx.opc.constants import CONTENT_TYPE as CT
from pptx.oxml import register_element_cls
from pptx import content_type_to_part_class_map
from pptx.opc.package import PartFactory

# path to media root for libreoffice to find files
MEDIA_ROOT = settings.MEDIA_ROOT

#global exclude to exclude punctuation
exclude = set(string.punctuation)

#for email
from django.core.mail import send_mail
from django.contrib.auth import models


class CT_NotesSlide(BaseOxmlElement):
	"""
	``<p:notes>`` element, root of a notesSlide part
	"""
	cSld = OneAndOnlyOne('p:cSld')
	clrMapOvr = ZeroOrOne('p:clrMapOvr', successors=(
		'p:transition', 'p:timing', 'p:extLst'
	))

	@classmethod
	def new(cls):
		"""
		Return a new ``<p:notes>`` element configured as a base slide shape.
		"""
		return parse_xml(cls._notes_xml())

	@staticmethod
	def _notes_xml():
		"""From http://msdn.microsoft.com/en-us/library/office/gg278319%28v=office.15%29.aspx#sectionSection4
		"""
		return (
		   '<p:notes %s>\n'
		   '   <p:cSld>\n'
		   '     <p:spTree>\n'
		   '       <p:nvGrpSpPr>\n'
		   '         <p:cNvPr id="1"\n'
		   '                  name="" />\n'
		   '         <p:cNvGrpSpPr />\n'
		   '         <p:nvPr />\n'
		   '       </p:nvGrpSpPr>\n'
		   '       <p:grpSpPr>\n'
		   '         <a:xfrm xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" />\n'
		   '       </p:grpSpPr>\n'
		   '       <p:sp>\n'
		   '         <p:nvSpPr>\n'
		   '           <p:cNvPr id="2"\n'
		   '                    name="" />\n'
		   '           <p:cNvSpPr>\n'
		   '             <a:spLocks noGrp="1"\n'
		   '                        xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" />\n'
		   '           </p:cNvSpPr>\n'
		   '           <p:nvPr>\n'
		   '             <p:ph />\n'
		   '           </p:nvPr>\n'
		   '         </p:nvSpPr>\n'
		   '         <p:spPr />\n'
		   '         <p:txBody>\n'
		   '           <a:bodyPr xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" />\n'
		   '           <a:lstStyle xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" />\n'
		   '           <a:p xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">\n'
		   '             <a:endParaRPr />\n'
		   '           </a:p>\n'
		   '         </p:txBody>\n'
		   '       </p:sp>\n'
		   '     </p:spTree>\n'
		   '   </p:cSld>\n'
		   '   <p:clrMapOvr>\n'
		   '     <a:masterClrMapping xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" />\n'
		   '  </p:clrMapOvr>\n'
		   '</p:notes>\n' % nsdecls('p', 'a', 'r')
		)


class SlideWrapper(Slide):
	"""Wrapper class for Slide that provides the .notes() property
	"""
	def __init__(self, slide):
		self._slide = slide
		self._notes_slide = None

	def __getattr__(self, attr):
		return self._slide.__getattribute__(attr)

	def notes_page(self):
		"""Return all related notesSlides
		"""
		notes_slide = None
		try:
			notes_slide = self.part_related_by(RT.NOTES_SLIDE)
		except KeyError:
			self.create_notes_slide()

		return self.part_related_by(RT.NOTES_SLIDE)


	def create_notes_slide(self):
		"""
		"""
		self._notes_slide = NotesSlide.new(self, self.partname, self.package)
		rId = self.relate_to(self._notes_slide, RT.NOTES_SLIDE)
		return rId


class NotesSlide(BaseSlide):
	"""This class will represent the Part of the notesSlide. Any notes retrieved
	from the presentation slides will be an instance of this class.
	"""
	@classmethod
	def new(cls, slide, partname, package):
		notes_slide_elm = CT_NotesSlide.new()
		slide = cls(partname, CT.PML_NOTES_SLIDE, notes_slide_elm, package)
		return slide

	@lazyproperty
	def shapes(self):
		"""
		Instance of |SlideShapeTree| containing sequence of shape objects
		appearing on this slide.
		"""
		return _SlideShapeTree(self)

	@lazyproperty
	def placeholders(self):
		"""
		Instance of |_SlidePlaceholders| containing sequence of placeholder
		shapes in this slide.
		"""
		return _SlidePlaceholders(self)

	def add_multiline_note(self, text):
		for line in text.split('\n'):
			self.add_note(line)

	def add_note(self, text):
		"""Add some text to the notesSlide, return paragraph
		that was added or False if no textframes were found
		"""
		for shape in self.shapes:
			if shape.has_text_frame and shape.is_placeholder:
				if hasattr(shape, 'ph_type') and shape.ph_type == PP_PLACEHOLDER_TYPE.BODY:
					return shape.text_frame.add_paragraph()
		return False

	def clear_notes(self):
		"""Remove all current notes from the slide
		"""
		for shape in self.shapes:
			if shape.has_text_frame:
				shape.text_frame.clear()

	def get_slide_runs(self):
		for shape in self.shapes:
			if shape.has_text_frame and shape.is_placeholder:
				for p in shape.text_frame.paragraphs:
					for run in p.runs:
						yield run


# Register the <p:notes> root tag
register_element_cls('p:notes', CT_NotesSlide)

# Add our NotesSlide as a valid Part for the notes Content Type
content_type_to_part_class_map[CT.PML_NOTES_SLIDE] = NotesSlide
PartFactory.part_type_for.update(content_type_to_part_class_map)


def strip_non_ascii(string):
	''' Returns the string without non ASCII characters'''
	stripped = (c for c in string if ord(c) >= 32 or ord(c) <= 126)
	return ''.join(stripped)



class Command(BaseCommand):

	def extract_text_images_from_pptx(self):
		prs = Presentation(self.lecturesObject.presentation)

		# convert lecturesObject.presentation to a pdf file in memory for imagemagick
		#where is the presentation?
		path_to_file = MEDIA_ROOT + '/' + str(self.lecturesObject.presentation)
		head, tail = os.path.split(path_to_file)
		print tail
		print path_to_file


		#name of pdf
		croppedtail = tail[:-5]
		pdffilename = croppedtail + ".pdf"

		# run the command and put the output in "head"
		# libreoffice --headless --convert-to pdf --outdir head path_to_file
		call(["soffice","--headless","--convert-to","pdf","--outdir",head,path_to_file])

		#extract all slides in pdf in imagemagic
		# convert head/pdffilename[index] head/slide.png
		pdfpath = head + '/' + pdffilename
		pngpath = head + "/slide.jpg"

		# only run the rest of script if pdf is present
		if os.path.isfile(pdfpath):
			call(["convert",pdfpath,pngpath])

			# text_runs will be populated with a list of strings,
			# one for each text run in presentation
			main_text = []
			notes_text = []
			slide_main_text = ''
			slide_notes_text = ''

			for index, slide in enumerate(prs.slides):
				print index
				# array for this slides text and notes to add the lectureSlides Object
				main_text_this_slide = []
				notes_text_this_slide = []
				# extract the text for each slide and append to array of text to save back to lecturesObject
				for shape in slide.shapes:
					if not shape.has_text_frame:
						continue
					for paragraph in shape.text_frame.paragraphs:
						for run in paragraph.runs:
							slide_main_text = strip_non_ascii(run.text)
							main_text_this_slide.append(slide_main_text)
							main_text.append(slide_main_text)

				notes_slide = SlideWrapper(slide).notes_page()
				for shape in notes_slide.shapes.__iter__():
					try:
						for paragraph in shape.text_frame.paragraphs:
							for run in paragraph.runs:
								slide_notes_text = strip_non_ascii(run.text)
								notes_text_this_slide.append(slide_notes_text)
								notes_text.append(slide_notes_text)
					except Exception:
						pass


				# add slide to lectureSlides
				addslide = lectureSlides()
				addslide.lecture = self.lecturesObject
				addslide.slide_number = index
				addslide.slide_main_text = '\n'.join(main_text_this_slide)
				addslide.slide_notes = '\n'.join(notes_text_this_slide)
				addslide.save()

				f = open(head + "/slide-"+ str(index) +".jpg")
				image_file = File(f)
				addslide.slide.save("slide-"+ str(index) +".jpg", image_file)
				f.close()
				os.remove(head + "/slide-"+ str(index) +".jpg")

				# for each slide create a pptx file with a single slide
				prs_slide = Presentation(self.lecturesObject.presentation)

				# calculate the last slide number
				largestslidenumber = len(prs_slide.slides)

				# loop through a range from the index to the largest slide number
				slide_number_plus_1 = index + 1
				for count in xrange(slide_number_plus_1, largestslidenumber):
					prs_slide.slides.delete_slide(prs_slide, slide_number_plus_1)

				for count in xrange(0, index):
					prs_slide.slides.delete_slide(prs_slide, 0)

				path_to_file = MEDIA_ROOT + '/' + str(croppedtail) + "_" + str(index) + ".pptx"
				prs_slide.save(path_to_file)

				# set prs_slide to None to free up memory
				prs_slide = None

				# run the garbage collector to free up memory
				gc.collect()

				seg = open(path_to_file)
				seg_file = File(seg)
				head_slide, tail_slide = os.path.split(path_to_file)
				print tail_slide
				addslide.presentation.save(tail_slide, seg_file)
				seg.close()
				os.remove(path_to_file)
				addslide.extracted = True
				addslide.save()


			#for text searching
			main_string = '\n'.join(main_text)
			notes_string = '\n'.join(notes_text)

			# concatonate strings together
			string = main_string + notes_string

			# save this string
			self.lecturesObject.presentation_text = string
			self.lecturesObject.extracted = True
			self.lecturesObject.save()

			# remove the large pdf to preserve server space
			os.remove(pdfpath)



	def extract_pptx_lectures(self):
		# pull PPTX from database where content hasn't yet been extracted
		lecturesObjects = lectures.objects.filter(extracted=False)

		# loop through lectures and pull all text
		for lecturesObject in lecturesObjects:
			try:
				self.lecturesObject = lecturesObject
				self.extract_text_images_from_pptx()
			except Exception as e:
				# email admins with an alert of the problem
				group = models.Group.objects.get(name="superusers")
				users = group.user_set.all()
				useremails = []
				for u in users:
					useremails.append(u.email)

				subject = "[GAHTC] There was an error importing a PPTX file."
				html_message = "Hello!<br /><br />There was an error importing the PPTX file located here:<br /><br />"+ str(lecturesObject.presentation) + "<br /><br />Please log in to the <a href='http://gahtc.org/dashboard/'>admin area</a> and review this file for any obvious problems. These may include files that are not .pptx files, or empty presentations. If you're unable to diagnose the problem, please contact JD Godchaux at jd@nijel.org.<br /><br />Thank you!<br />GAHTC | Global Architectural History Teaching Collaborative<br /><br />Error message:" + str(e)
				message = "Hello! There was an error importing the PPTX file located here: "+ str(lecturesObject.presentation) + ". Please log in to the admin area: http://gahtc.org/dashboard/ and review this file for any obvious problems. These may include files that are not .pptx files, or empty presentations. If you're unable to diagnose the problem, please contact JD Godchaux at jd@nijel.org. Thank you! GAHTC | Global Architectural History Teaching Collaborative. Error message:" + str(e)

				print e

				send_mail(subject, message, 'gahtcweb@gmail.com', useremails, fail_silently=True, html_message=html_message)

				pass


	def handle(self, *args, **options):
		print "Extracting Lecture PPTX files...."
		self.extract_pptx_lectures()
		print "Done."



