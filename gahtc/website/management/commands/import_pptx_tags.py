import sys,os
from django.core.management.base import BaseCommand, CommandError
from django.conf import settings
from django.core.files import File
from subprocess import call
from website.models import *
import pptx
from pptx import Presentation
import string
from textblob import TextBlob

"""
  Opens pptx files and extracts all text
  Requires libreoffice and Imagemagick
"""

from pptx.util import lazyproperty, Pt
from pptx.parts.slide import BaseSlide, Slide, _SlideShapeTree, _SlidePlaceholders
from pptx.shapes.shape import BaseShape
from pptx.opc.constants import RELATIONSHIP_TYPE as RT, CONTENT_TYPE as CT
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER_TYPE
from pptx.oxml.xmlchemy import BaseOxmlElement, OneAndOnlyOne, ZeroOrOne
from pptx.oxml import parse_xml
from pptx.oxml.ns import nsdecls
from pptx.opc.constants import CONTENT_TYPE as CT
from pptx.oxml import register_element_cls
from pptx import content_type_to_part_class_map
from pptx.opc.package import PartFactory

# path to media root for libreoffice to find files
MEDIA_ROOT = settings.MEDIA_ROOT

#global exclude to exclude punctuation
exclude = set(string.punctuation)


class CT_NotesSlide(BaseOxmlElement):
    """
    ``<p:notes>`` element, root of a notesSlide part
    """
    cSld = OneAndOnlyOne('p:cSld')
    clrMapOvr = ZeroOrOne('p:clrMapOvr', successors=(
        'p:transition', 'p:timing', 'p:extLst'
    ))

    @classmethod
    def new(cls):
        """
        Return a new ``<p:notes>`` element configured as a base slide shape.
        """
        return parse_xml(cls._notes_xml())

    @staticmethod
    def _notes_xml():
        """From http://msdn.microsoft.com/en-us/library/office/gg278319%28v=office.15%29.aspx#sectionSection4
        """
        return (
           '<p:notes %s>\n'
           '   <p:cSld>\n'
           '     <p:spTree>\n'
           '       <p:nvGrpSpPr>\n'
           '         <p:cNvPr id="1"\n'
           '                  name="" />\n'
           '         <p:cNvGrpSpPr />\n'
           '         <p:nvPr />\n'
           '       </p:nvGrpSpPr>\n'
           '       <p:grpSpPr>\n'
           '         <a:xfrm xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" />\n'
           '       </p:grpSpPr>\n'
           '       <p:sp>\n'
           '         <p:nvSpPr>\n'
           '           <p:cNvPr id="2"\n'
           '                    name="" />\n'
           '           <p:cNvSpPr>\n'
           '             <a:spLocks noGrp="1"\n'
           '                        xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" />\n'
           '           </p:cNvSpPr>\n'
           '           <p:nvPr>\n'
           '             <p:ph />\n'
           '           </p:nvPr>\n'
           '         </p:nvSpPr>\n'
           '         <p:spPr />\n'
           '         <p:txBody>\n'
           '           <a:bodyPr xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" />\n'
           '           <a:lstStyle xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" />\n'
           '           <a:p xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">\n'
           '             <a:endParaRPr />\n'
           '           </a:p>\n'
           '         </p:txBody>\n'
           '       </p:sp>\n'
           '     </p:spTree>\n'
           '   </p:cSld>\n'
           '   <p:clrMapOvr>\n'
           '     <a:masterClrMapping xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" />\n'
           '  </p:clrMapOvr>\n'
           '</p:notes>\n' % nsdecls('p', 'a', 'r')
        )


class SlideWrapper(Slide):
    """Wrapper class for Slide that provides the .notes() property
    """
    def __init__(self, slide):
        self._slide = slide
        self._notes_slide = None

    def __getattr__(self, attr):
        return self._slide.__getattribute__(attr)

    def notes_page(self):
        """Return all related notesSlides
        """
        notes_slide = None
        try:
            notes_slide = self.part_related_by(RT.NOTES_SLIDE)
        except KeyError:
            self.create_notes_slide()

        return self.part_related_by(RT.NOTES_SLIDE)


    def create_notes_slide(self):
        """
        """
        self._notes_slide = NotesSlide.new(self, self.partname, self.package)
        rId = self.relate_to(self._notes_slide, RT.NOTES_SLIDE)
        return rId


class NotesSlide(BaseSlide):
    """This class will represent the Part of the notesSlide. Any notes retrieved
    from the presentation slides will be an instance of this class.
    """
    @classmethod
    def new(cls, slide, partname, package):
        notes_slide_elm = CT_NotesSlide.new()
        slide = cls(partname, CT.PML_NOTES_SLIDE, notes_slide_elm, package)
        return slide

    @lazyproperty
    def shapes(self):
        """
        Instance of |_SlideShapeTree| containing sequence of shape objects
        appearing on this slide.
        """
        return _SlideShapeTree(self)

    @lazyproperty
    def placeholders(self):
        """
        Instance of |_SlidePlaceholders| containing sequence of placeholder
        shapes in this slide.
        """
        return _SlidePlaceholders(self)

    def add_multiline_note(self, text):
        for line in text.split('\n'):
            self.add_note(line)

    def add_note(self, text):
        """Add some text to the notesSlide, return paragraph
        that was added or False if no textframes were found
        """
        for shape in self.shapes:
            if shape.has_text_frame and shape.is_placeholder:
                if hasattr(shape, 'ph_type') and shape.ph_type == PP_PLACEHOLDER_TYPE.BODY:
                    return shape.text_frame.add_paragraph()
        return False

    def clear_notes(self):
        """Remove all current notes from the slide
        """
        for shape in self.shapes:
            if shape.has_text_frame:
                shape.text_frame.clear()

    def get_slide_runs(self):
        for shape in self.shapes:
            if shape.has_text_frame and shape.is_placeholder:
                for p in shape.text_frame.paragraphs:
                    for run in p.runs:
                        yield run


# Register the <p:notes> root tag
register_element_cls('p:notes', CT_NotesSlide)

# Add our NotesSlide as a valid Part for the notes Content Type
content_type_to_part_class_map[CT.PML_NOTES_SLIDE] = NotesSlide
PartFactory.part_type_for.update(content_type_to_part_class_map)


def strip_non_ascii(string):
    ''' Returns the string without non ASCII characters'''
    stripped = (c for c in string if ord(c) >= 32 or ord(c) <= 126)
    return ''.join(stripped)



class Command(BaseCommand):
    
    def extract_tags_from_pptx_lectures(self):
        # pull PPTX from database where content hasn't yet been extracted
        lecturesObjects = lectures.objects.all()

        # loop through lectures and pull all text
        for lecturesObject in lecturesObjects:
            if not lecturesObject.tags:
                prs = Presentation(lecturesObject.presentation)

                # text_runs will be populated with a list of strings,
                # one for each text run in presentation
                main_text = []
                notes_text = []
                slide_main_text = ''
                slide_notes_text = ''

                for index, slide in enumerate(prs.slides):
                    # array for this slides text and notes to add the lectureSlides Object
                    main_text_this_slide = []
                    notes_text_this_slide = []
                    # extract the text for each slide and append to array of text to save back to lecturesObject
                    for shape in slide.shapes:
                        if not shape.has_text_frame:
                            continue
                        for paragraph in shape.text_frame.paragraphs:
                            for run in paragraph.runs:
                                slide_main_text = strip_non_ascii(run.text)
                                main_text_this_slide.append(slide_main_text)
                                main_text.append(slide_main_text)

                    notes_slide = SlideWrapper(slide).notes_page()
                    for shape in notes_slide.shapes.__iter__():
                        for paragraph in shape.text_frame.paragraphs:
                            for run in paragraph.runs:
                                slide_notes_text = strip_non_ascii(run.text)
                                notes_text_this_slide.append(slide_notes_text)
                                notes_text.append(slide_notes_text)

                    # add slide to lectureSlides
                    print lecturesObject
                    print index
                    addslide = lectureSlides.objects.get(lecture=lecturesObject, slide_number=index)
                    # add in tags for this slide
                    # create tags from noun_phrases for main slide text
                    mt = '. '.join(main_text_this_slide)
                    blobbed = TextBlob(mt)
                    np = blobbed.noun_phrases
                    np = list(set(np))
                    np = [s for s in np if s]
                    addslide.tags.clear()
                    for item in np:
                        s = ''.join(ch for ch in item if ch not in exclude)
                        addslide.tags.add(s)

                    # create tags from noun_phrases for slide notes
                    nt = '. '.join(notes_text_this_slide)
                    blobbed = TextBlob(nt)
                    np = blobbed.noun_phrases
                    np = list(set(np))
                    np = [s for s in np if s]
                    for item in np:
                        s = ''.join(ch for ch in item if ch not in exclude)
                        addslide.tags.add(s)


                #add tags for entire presentation
                # create tags from noun_phrases for main slide text
                mt = '. '.join(main_text)
                blobbed = TextBlob(mt)
                np = blobbed.noun_phrases
                np = list(set(np))
                np = [s for s in np if s]
                lecturesObject.tags.clear()
                for item in np:
                    s = ''.join(ch for ch in item if ch not in exclude)
                    lecturesObject.tags.add(s)

                # create tags from noun_phrases for slide notes
                nt = '. '.join(notes_text)
                blobbed = TextBlob(nt)
                np = blobbed.noun_phrases
                np = list(set(np))
                np = [s for s in np if s]
                for item in np:
                    s = ''.join(ch for ch in item if ch not in exclude)
                    lecturesObject.tags.add(s)               


    def extract_tags_from_pptx_segments(self):
        # pull PPTX from database where content hasn't yet been extracted
        lecturesObjects = lectureSegments.objects.all()

        # loop through lectures and pull all text
        for lecturesObject in lecturesObjects:
            if not lecturesObject.tags:
                prs = Presentation(lecturesObject.presentation)

                # text_runs will be populated with a list of strings,
                # one for each text run in presentation
                main_text = []
                notes_text = []
                slide_main_text = ''
                slide_notes_text = ''

                for index, slide in enumerate(prs.slides):
                    # array for this slides text and notes to add the lectureSlides Object
                    main_text_this_slide = []
                    notes_text_this_slide = []
                    # extract the text for each slide and append to array of text to save back to lecturesObject
                    for shape in slide.shapes:
                        if not shape.has_text_frame:
                            continue
                        for paragraph in shape.text_frame.paragraphs:
                            for run in paragraph.runs:
                                slide_main_text = strip_non_ascii(run.text)
                                main_text_this_slide.append(slide_main_text)
                                main_text.append(slide_main_text)

                    notes_slide = SlideWrapper(slide).notes_page()
                    for shape in notes_slide.shapes.__iter__():
                        for paragraph in shape.text_frame.paragraphs:
                            for run in paragraph.runs:
                                slide_notes_text = strip_non_ascii(run.text)
                                notes_text_this_slide.append(slide_notes_text)
                                notes_text.append(slide_notes_text)


                #add tags for entire presentation
                # create tags from noun_phrases for main slide text
                mt = '. '.join(main_text)
                blobbed = TextBlob(mt)
                np = blobbed.noun_phrases
                np = list(set(np))
                np = [s for s in np if s]
                lecturesObject.tags.clear()
                for item in np:
                    s = ''.join(ch for ch in item if ch not in exclude)
                    lecturesObject.tags.add(s)

                # create tags from noun_phrases for slide notes
                nt = '. '.join(notes_text)
                blobbed = TextBlob(nt)
                np = blobbed.noun_phrases
                np = list(set(np))
                np = [s for s in np if s]
                for item in np:
                    s = ''.join(ch for ch in item if ch not in exclude)
                    lecturesObject.tags.add(s)               


    def handle(self, *args, **options):
        print "Extracting Tags for Lecture PPTX files...."
        self.extract_tags_from_pptx_lectures()
        print "Extracting Tags for Lecture Segment PPTX files...."
        self.extract_tags_from_pptx_segments()
        print "Done."



